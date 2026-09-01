# -*- coding: utf-8 -*-
"""内部知识库 Flask Blueprint。"""

import hashlib
import io
import json
import math
import os
import re
import secrets
import shutil
import struct
import subprocess
import tempfile
import threading
import time
import unicodedata
import urllib.request
import urllib.error
import uuid
from collections import Counter
from datetime import datetime, timezone, timedelta
from pathlib import Path

from flask import (Blueprint, current_app, request, send_file, send_from_directory,
                   jsonify, session, redirect, url_for, abort, Response,
                   stream_with_context)
from werkzeug.security import check_password_hash, generate_password_hash
from openpyxl import Workbook
from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
from openpyxl.utils import get_column_letter

from paths import (INTERNAL_KNOWLEDGE_BASE_DIR, INTERNAL_KNOWLEDGE_BASE_DB,
                   INTERNAL_KNOWLEDGE_BASE_UPLOADS, INTERNAL_KNOWLEDGE_BASE_PDF_CACHE,
                   INTERNAL_KNOWLEDGE_BASE_TEMP)
from .storage import SQLiteStore

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
DATA_DIR = str(INTERNAL_KNOWLEDGE_BASE_DIR)
CACHE_DIR = str(INTERNAL_KNOWLEDGE_BASE_TEMP)
PREVIEW_CACHE_DIR = str(INTERNAL_KNOWLEDGE_BASE_PDF_CACHE)
UPLOAD_DIR = str(INTERNAL_KNOWLEDGE_BASE_UPLOADS)
STORE_PATH = str(INTERNAL_KNOWLEDGE_BASE_DB)
CONVERSION_VERSION = "libreoffice-v1"

CST = timezone(timedelta(hours=8))

# 可转换的文件扩展名
CONVERTIBLE_EXTS = {".pptx", ".ppt", ".doc", ".docx", ".xls", ".xlsx"}
ALLOWED_UPLOAD_EXTS = {".pdf", ".ppt", ".pptx", ".doc", ".docx", ".xls", ".xlsx"}
MAX_UPLOAD_SIZE = 100 * 1024 * 1024
DEFAULT_USER_PASSWORD = "123456"

# 报告主题分类（与前端 js/data.js 中 REPORT_THEMES 保持一致）
REPORT_THEMES = (
    "macro_rate", "credit", "multi_asset", "quant", "fixed_income_plus",
    "equity", "other",
)
REPORT_THEME_LABELS = {
    "macro_rate": "宏观利率", "credit": "信用", "multi_asset": "多资产", "quant": "量化",
    "fixed_income_plus": "固收+", "equity": "权益", "other": "其他",
}
REPORT_CATEGORY_LABELS = {
    "weekly": "周报", "monthly": "月报", "deep": "深度报告", "other": "其他报告",
}


def _default_reminder_config():
    """附件《报告要求.xlsx》对应的默认专题报告规则。"""
    return {
        "period": datetime.now(CST).strftime("%Y"),
        "reportCategory": "deep",
        "rules": [
            {"id": "fi-credit", "label": "固收中心 · 信用组", "mode": "group", "target": 2,
             "userIds": ["huyongyan", "zhenghongbin"]},
            {"id": "fi-multi-asset", "label": "固收中心 · 多资产组", "mode": "group", "target": 2,
             "userIds": ["qianyouni", "cuizhuoju", "zhangyunfan"]},
            {"id": "fi-quant", "label": "固收中心 · 量化组", "mode": "group", "target": 2,
             "userIds": ["zhengmeng", "maxiao", "ouyangquan"]},
            {"id": "fi-plus", "label": "固收中心 · 固收+组", "mode": "group", "target": 2,
             "userIds": ["maying", "linyuang"]},
            {"id": "fi-macro-rate", "label": "固收中心 · 宏观利率组", "mode": "group", "target": 2,
             "userIds": ["kanghairong", "zhumeng"]},
            {"id": "aa-feilixiao", "label": "资产配置部研究组 · 费立孝", "mode": "person", "target": 2,
             "userIds": ["feilixiao"]},
            {"id": "aa-zhangqingchang", "label": "资产配置部研究组 · 张庆昌", "mode": "person", "target": 2,
             "userIds": ["zhangqingchang"]},
            {"id": "aa-wanghui", "label": "资产配置部研究组 · 王辉", "mode": "person", "target": 2,
             "userIds": ["wanghui"]},
            {"id": "aa-liushengyao", "label": "资产配置部研究组 · 刘圣尧", "mode": "person", "target": 2,
             "userIds": ["liushengyao"]},
            {"id": "aa-zongshaohui", "label": "资产配置部研究组 · 宗韶晖", "mode": "person", "target": 2,
             "userIds": ["zongshaohui"]},
            {"id": "aa-songlingfeng", "label": "资产配置部研究组 · 宋凌峰", "mode": "person", "target": 2,
             "userIds": ["songlingfeng"]},
        ],
    }

os.makedirs(CACHE_DIR, exist_ok=True)
os.makedirs(PREVIEW_CACHE_DIR, exist_ok=True)
os.makedirs(UPLOAD_DIR, exist_ok=True)


# --------------------------------------------------------------------------- #
# 大模型配置统一在项目根 llm_config.py：四个模型（自部署 / DeepSeek 内网部署 /
# MiMo / DeepSeek 官方），优先级由后台"大模型管理"持久化，调用时实时读取。
# --------------------------------------------------------------------------- #
import llm_config


def _llm_api_key():
    """任一模型密钥可用即视为大模型已配置。"""
    return llm_config.llm_api_key()


LLM_TIMEOUT = 90
LLM_MAX_TEXT_CHARS = 6000  # 发送给 LLM 的文档文本最大长度
KNOWLEDGE_INDEX_TEXT_CHARS = 60000
KNOWLEDGE_VECTOR_DIM = 512
KNOWLEDGE_VECTOR_VERSION = "local-char-ngram-v1-d512"
KNOWLEDGE_CHUNK_CHARS = 1800
KNOWLEDGE_CHUNK_OVERLAP = 240

# 单篇报告 AI 摘要：三个篇幅版本 + 送入模型的全文上限
AI_SUMMARY_TEXT_CHARS = 60000
AI_SUMMARY_STYLES = {
    "concise": {"label": "精炼版", "max_tokens": 1200},
    "standard": {"label": "标准版", "max_tokens": 2200},
    "deep": {"label": "深度版", "max_tokens": 3200},
}
REPORT_TYPE_LABELS = {
    "internal": "内部报告", "external": "外部报告",
    "research_visit": "调研报告", "roadshow": "路演报告",
}
AI_SUMMARY_SYSTEM = (
    "你是固收投研团队的研究主管，为团队成员快速消化研究报告撰写详细摘要。"
    "严格忠于原文，不编造数据、观点或结论，关键处保留具体数字。"
)
AI_SUMMARY_STYLE_PROMPTS = {
    "concise": (
        "请输出结构化中文摘要（markdown 格式），仅包含以下两节：\n"
        "### 核心结论\n（3-5 条，每条一句话给出最重要的结论或判断）\n"
        "### 关键数据\n（支撑结论的关键数据与事实，保留具体数字）\n"
        "总篇幅控制在 400-600 字。"
    ),
    "standard": (
        "请输出结构化中文摘要（markdown 格式），包含以下四节：\n"
        "### 核心结论\n（3-5 条，每条一句话给出最重要的结论或判断）\n"
        "### 主要观点与分析逻辑\n（按报告行文顺序提炼主要论点、分析框架与推理链条）\n"
        "### 关键数据与论据\n（支撑结论的关键数据、事实与证据，保留具体数字）\n"
        "### 风险提示与关注点\n（报告提及或隐含的风险、假设与局限）\n"
        "总篇幅控制在 600-1000 字。"
    ),
    "deep": (
        "请输出结构化中文摘要（markdown 格式），包含以下四节：\n"
        "### 核心结论\n（3-5 条，每条一句话给出最重要的结论或判断）\n"
        "### 分章节要点\n（按报告章节逐段提炼，尽量覆盖全文内容）\n"
        "### 关键数据与论据\n（支撑结论的关键数据、事实与证据，保留具体数字）\n"
        "### 风险提示与关注点\n（报告提及或隐含的风险、假设与局限）\n"
        "总篇幅控制在 1000-1500 字。"
    ),
}
AI_SUMMARY_COMMON_RULES = (
    "其他要求：\n"
    "- 调研/路演纪要类报告，重点提炼交流对象的核心观点与问答要点\n"
    "- 全文过长被截断时，优先覆盖核心章节\n"
    "- 用中文，严格忠于原文"
)
_knowledge_index_lock = threading.Lock()


# --------------------------------------------------------------------------- #
# LibreOffice 查找与转换
# --------------------------------------------------------------------------- #
def find_soffice():
    # LIBREOFFICE_PATH 显式指定（.env / 服务环境），优先级最高；
    # 其后依次尝试标准安装目录与免安装解压目录（D:\LibreOffice），最后搜 PATH。
    explicit = os.environ.get("LIBREOFFICE_PATH", "").strip()
    candidates = [explicit] if explicit else []
    candidates += [
        r"C:\Program Files\LibreOffice\program\soffice.com",
        r"C:\Program Files\LibreOffice\program\soffice.exe",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.com",
        r"C:\Program Files (x86)\LibreOffice\program\soffice.exe",
        r"D:\LibreOffice\program\soffice.com",
        r"D:\LibreOffice\program\soffice.exe",
        r"D:\LibreOffice\LibreOffice\program\soffice.com",
        r"D:\LibreOffice\LibreOffice\program\soffice.exe",
    ]
    for path in candidates:
        if os.path.isfile(path):
            return path
    for name in ("soffice", "soffice.com", "soffice.exe", "libreoffice"):
        found = shutil.which(name)
        if found:
            return found
    return None


SOFFICE = find_soffice()


def source_hash(file_path):
    digest = hashlib.sha256()
    with open(file_path, "rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def cache_key(file_path, file_size=None):
    digest = source_hash(file_path)
    return hashlib.sha256(f"{CONVERSION_VERSION}:{digest}".encode("utf-8")).hexdigest()


def convert_to_pdf(input_path, output_dir):
    # 每次转换使用独立 LibreOffice 用户配置，避免桌面 Office/并发预览导致
    # soffice 复用或锁定全局配置后无输出退出。配置目录随预览临时目录清理。
    profile_dir = os.path.join(output_dir, ".lo_profile")
    os.makedirs(profile_dir, exist_ok=True)
    profile_uri = Path(profile_dir).resolve().as_uri()
    cmd = [SOFFICE, f"-env:UserInstallation={profile_uri}",
           "--headless", "--norestore", "--nolockcheck",
           "--convert-to", "pdf", "--outdir", output_dir, input_path]
    # soffice 在中文 Windows 上按控制台代码页（GBK）输出日志，显式容错解码，
    # 避免读取 stderr 的线程因编码异常崩溃。
    result = subprocess.run(
        cmd, capture_output=True, timeout=120,
        encoding="utf-8", errors="replace",
    )
    if result.returncode != 0:
        raise RuntimeError(f"LibreOffice 转换失败: {result.stderr or result.stdout}")
    base_name = os.path.splitext(os.path.basename(input_path))[0]
    pdf_path = os.path.join(output_dir, base_name + ".pdf")
    if not os.path.isfile(pdf_path):
        raise RuntimeError("转换后未找到 PDF 文件")
    return pdf_path


def _remove_preview_artifacts(path, retries=30, retry_delay=1):
    """删除单次预览产生的临时目录；Windows 文件句柄未释放时自动重试。"""
    for _ in range(retries):
        try:
            if os.path.isdir(path):
                shutil.rmtree(path)
            elif os.path.exists(path):
                os.remove(path)
            return True
        except OSError:
            time.sleep(retry_delay)
    return False


def _schedule_preview_cleanup(path):
    """响应完成即清理，并提供后台重试兜底，最迟约 30 秒再次尝试。"""
    worker = threading.Thread(target=_remove_preview_artifacts, args=(path,), daemon=True)
    worker.start()


def _send_temporary_pdf(pdf_path, work_dir):
    response = send_file(pdf_path, mimetype="application/pdf", as_attachment=False)
    response.call_on_close(lambda: _remove_preview_artifacts(work_dir, retries=1, retry_delay=0))
    _schedule_preview_cleanup(work_dir)
    return response


def _make_preview_work_dir():
    work_dir = os.path.join(CACHE_DIR, f"preview_{uuid.uuid4().hex}")
    os.makedirs(work_dir, exist_ok=False)
    return work_dir


def clear_preview_cache():
    """启动时清除上次异常退出可能遗留的一次性预览临时目录。

    正式 PDF 缓存位于独立目录，不参与启动清理并永久保留，直至超级管理员删除。
    """
    removed = 0
    if not os.path.isdir(CACHE_DIR):
        return removed
    for name in os.listdir(CACHE_DIR):
        if name == os.path.basename(PREVIEW_CACHE_DIR):
            continue  # 跳过持久化 PDF 缓存目录
        path = os.path.join(CACHE_DIR, name)
        try:
            if os.path.isdir(path):
                shutil.rmtree(path)
            else:
                os.remove(path)
            removed += 1
        except OSError:
            pass
    return removed


def _get_cached_pdf(file_path):
    """返回永久缓存中的 PDF；缓存不按时间自动过期。"""
    key = cache_key(file_path)
    cached = os.path.join(PREVIEW_CACHE_DIR, f"{key}.pdf")
    if not os.path.isfile(cached):
        if store.get_pdf_cache(key):
            store.delete_pdf_cache(key)
        return None
    store.touch_pdf_cache(key)
    return cached


def _store_cached_pdf(src_pdf, file_path, report_id=None):
    """将转换产物落盘到持久缓存，返回缓存路径；失败时返回 None 不影响预览。"""
    key = cache_key(file_path)
    cached = os.path.join(PREVIEW_CACHE_DIR, f"{key}.pdf")
    try:
        # 同目录 move 最快；跨目录先 copy 再清理源
        if os.path.abspath(os.path.dirname(src_pdf)) == os.path.abspath(PREVIEW_CACHE_DIR):
            cached = src_pdf
        else:
            shutil.move(src_pdf, cached)
        store.upsert_pdf_cache(
            key, report_id, source_hash(file_path), os.path.basename(cached),
            os.path.getsize(cached), CONVERSION_VERSION,
        )
        return cached
    except OSError:
        return None


def _preview_cache_cleanup_loop():
    """兼容旧调用；正式 PDF 缓存只允许超级管理员删除。"""
    return None


def _delete_cache_files(cache_keys):
    deleted = 0
    root = Path(PREVIEW_CACHE_DIR).resolve()
    for key in cache_keys:
        if not re.fullmatch(r"[0-9a-f]{64}", str(key)):
            continue
        path = (root / f"{key}.pdf").resolve()
        if root not in path.parents:
            continue
        try:
            path.unlink(missing_ok=True)
            deleted += 1
        except OSError:
            current_app.logger.exception("删除 PDF 缓存失败: %s", key)
    return deleted


# --------------------------------------------------------------------------- #
# 数据持久化
# --------------------------------------------------------------------------- #
def _now_iso():
    return datetime.now(CST).isoformat(timespec="seconds")


def _infer_theme(report):
    """根据报告标题、摘要、标签推断主题分类（用于迁移旧数据）。"""
    text = " ".join([
        report.get("title", ""),
        report.get("summary", ""),
        " ".join(report.get("tags") or []),
    ]).lower()
    if any(k in text for k in ["信用", "credit", "利差"]):
        return "credit"
    if any(k in text for k in ["固收+", "固收加", "fixed_income_plus"]):
        return "fixed_income_plus"
    if any(k in text for k in ["多资产", "multi_asset", "股债"]):
        return "multi_asset"
    if any(k in text for k in ["量化", "quant", "模型"]):
        return "quant"
    if any(k in text for k in ["权益", "equity", "股票", "行业"]):
        return "equity"
    if any(k in text for k in ["宏观", "利率", "macro", "经济周期", "利率债"]):
        return "macro_rate"
    return "other"


clear_preview_cache()
store = SQLiteStore(Path(STORE_PATH), _default_reminder_config())
bp = Blueprint("internal_knowledge_base", __name__)
app = bp  # Keep existing route declarations compact.


def csrf_token():
    token = session.get("internal_knowledge_base_csrf")
    if not token:
        token = secrets.token_urlsafe(32)
        session["internal_knowledge_base_csrf"] = token
    return token


@bp.before_request
def protect_internal_knowledge_base():
    if not session.get("authenticated"):
        if "/api/" in request.path:
            return jsonify({"error": "请先登录内部研究工作台", "portalAuthRequired": True}), 401
        return redirect(url_for("login", next=request.full_path.rstrip("?")))
    if request.content_length and request.content_length > MAX_UPLOAD_SIZE + 1024 * 1024:
        return jsonify({"error": "上传文件超过100MB限制"}), 413
    if request.method not in {"GET", "HEAD", "OPTIONS"} and "/api/" in request.path:
        submitted = request.headers.get("X-CSRF-Token", "")
        expected = session.get("internal_knowledge_base_csrf", "")
        if not expected or not secrets.compare_digest(submitted, expected):
            return jsonify({"error": "安全令牌已失效，请刷新页面后重试"}), 403


@bp.after_request
def audit_mutation(response):
    if request.method not in {"GET", "HEAD", "OPTIONS"} and "/api/" in request.path and response.status_code < 400:
        actor_id = session.get("internal_knowledge_base_user_id")
        actor_type = "superadmin" if session.get("internal_knowledge_base_superadmin") else "user"
        try:
            store.audit(actor_type, actor_id, request.endpoint or request.path,
                        "request", next(iter((request.view_args or {}).values()), None),
                        {"method": request.method}, request.remote_addr)
        except Exception:
            current_app.logger.exception("内部知识库审计日志写入失败")
    return response

# --------------------------------------------------------------------------- #
# 辅助函数
# --------------------------------------------------------------------------- #
def public_user(user):
    """对外返回用户信息（不含密码）。"""
    if not user:
        return None
    return {k: user.get(k) for k in ("id", "name", "org", "role")}


def public_report(report):
    """对外返回报告信息（不含内部字段）。"""
    if not report:
        return None
    return dict(report)


SCORING_ORGS = ("资产配置部", "固收中心")


def report_scoring_orgs(report):
    """返回报告的打分部门列表；未配置时默认两部门均打分（向后兼容）。"""
    orgs = report.get("scoringOrgs")
    if not isinstance(orgs, list) or not orgs:
        return list(SCORING_ORGS)
    return [org for org in orgs if org in SCORING_ORGS] or list(SCORING_ORGS)


def leader_in_scoring_scope(user, scoring_orgs):
    """部门领导按报告所选打分部门过滤；不隶属具体打分部门的通用领导始终参与。"""
    org = user.get("org")
    if org in SCORING_ORGS:
        return org in scoring_orgs
    return True


def eligible_scorers(report):
    """返回某份月报/深度报告的应评分人员列表。

    规则：部门领导（org 为资产配置部/固收中心）仅当其所属部门在报告所选打分
    部门内时参与；通用领导（org 为领导/行政/空）始终参与；研究人员仅当其所属
    部门在报告所选打分部门内时参与。行政账号不参与评分。报告对所有人可见，
    本函数仅决定评分资格。
    """
    scoring_orgs = set(report_scoring_orgs(report))
    result = []
    for user in store.users():
        role = user.get("role")
        if role == "admin":
            continue
        if role == "leader":
            if leader_in_scoring_scope(user, scoring_orgs):
                result.append(user)
        elif user.get("org") in scoring_orgs:
            result.append(user)
    return result


def is_eligible_scorer(report, user):
    """判断指定用户是否具备该报告的评分资格。"""
    if not user or user.get("role") == "admin":
        return False
    if user.get("role") == "leader":
        return leader_in_scoring_scope(user, set(report_scoring_orgs(report)))
    return user.get("org") in set(report_scoring_orgs(report))


def report_file_path(report):
    """返回报告文件的安全绝对路径；文件无效时返回空串。"""
    rel = str(report.get("fileUrl", ""))
    if not rel:
        return ""
    upload_root = Path(UPLOAD_DIR).resolve()
    full = (upload_root / Path(rel).name).resolve()
    if upload_root not in full.parents:
        return ""
    return str(full) if full.is_file() else ""


def format_bytes(size):
    try:
        size = float(size)
    except (TypeError, ValueError):
        return "—"
    if size <= 0:
        return "0 KB"
    import math
    units = ["B", "KB", "MB", "GB"]
    idx = min(int(math.log(size, 1024)), len(units) - 1)
    val = size / (1024 ** idx)
    fmt = f"{val:.1f}" if idx > 1 else f"{val:.0f}"
    if fmt.endswith(".0"):
        fmt = fmt[:-2]
    return f"{fmt} {units[idx]}"


def require_user():
    """要求普通用户已登录，返回用户对象或 None。"""
    uid = session.get("internal_knowledge_base_user_id")
    if not uid:
        return None
    return store.get_user(uid)


def require_role(*roles):
    user = require_user()
    if not user or user.get("role") not in roles:
        return None
    return user


def require_admin():
    return True if session.get("internal_knowledge_base_superadmin") else None


def safe_filename(name):
    """把文件名清洗为安全的文件系统名（保留中文）。"""
    name = os.path.basename(name)
    name = re.sub(r'[\\/:*?"<>|\r\n\t]', "_", name)
    return name.strip().strip(".") or "upload"


def stored_filename(report_id, original_name):
    """生成上传文件在 uploads/ 下的保存名。"""
    return f"{report_id}__{safe_filename(original_name)}"


# --------------------------------------------------------------------------- #
# 文档文本抽取 & DeepSeek 智能补全
# --------------------------------------------------------------------------- #
def _extract_text(file_path, max_chars=LLM_MAX_TEXT_CHARS):
    """从 PDF/DOCX/PPTX/XLSX 中抽取纯文本，截断到 max_chars。
    抽取失败时返回空串，由调用方用文件名兜底。
    """
    ext = os.path.splitext(file_path)[1].lower()
    try:
        if ext == ".pdf":
            import pymupdf
            parts = []
            with pymupdf.open(file_path) as doc:
                for page in doc:
                    parts.append(page.get_text("text"))
                    if sum(len(p) for p in parts) > max_chars:
                        break
            text = "\n".join(parts)
        elif ext == ".docx":
            import docx
            doc = docx.Document(file_path)
            text = "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        elif ext in (".pptx", ".ppt"):
            import pptx
            prs = pptx.Presentation(file_path)
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        parts.append(shape.text_frame.text)
                if sum(len(p) for p in parts) > max_chars:
                    break
            text = "\n".join(parts)
        elif ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        parts.append("\t".join(cells))
                    if sum(len(p) for p in parts) > max_chars:
                        break
                if sum(len(p) for p in parts) > max_chars:
                    break
            wb.close()
            text = "\n".join(parts)
        else:
            text = ""
    except Exception:
        text = ""
    return text[:max_chars]


def _strip_code_fences(text):
    """部分网关（自部署 GLM）不支持 response_format，模型会把 JSON 包在
    ```json 围栏里输出；剥掉围栏再交给调用方 json.loads。"""
    s = str(text or "").strip()
    if s.startswith("```"):
        s = re.sub(r"^```[a-zA-Z]*\s*", "", s)
        if s.endswith("```"):
            s = s[:-3]
    return s.strip()


def _call_llm(prompt, system=None, max_tokens=None, json_mode=True):
    """调用大模型（自部署优先，MiMo 兜底，DeepSeek 二层兜底），返回 message content。

    system：可选的系统消息，用于约束角色与回答风格。
    max_tokens：可选的输出长度上限，未传时由模型默认值决定。
    json_mode：是否要求 JSON 输出；仅支持 response_format 的 provider 会带上
    该参数，其余 provider（自部署网关不支持 response_format）靠提示词约束，
    返回前统一剥离 ``` 围栏。知识问答需要流式展示纯文本，传 False。
    """
    messages = _llm_messages(prompt, system=system)
    providers = llm_config.available_providers()
    if not providers:
        raise RuntimeError("未配置任何大模型（SELF_LLM_API_KEY / DEEPSEEK_INTERNAL_BASE_URL / MIMO_API_KEY / DEEPSEEK_API_KEY）")
    last_error = None
    for provider in providers:
        try:
            content = _complete_provider(provider, messages, max_tokens, json_mode)
            return _strip_code_fences(content) if json_mode else content
        except Exception as exc:
            last_error = exc
            current_app.logger.warning("%s 调用失败，尝试下一个大模型：%s", provider["name"], exc)
    raise RuntimeError(f"大模型调用失败：{last_error}")


def _stream_llm(prompt, system=None, max_tokens=None, provider_sink=None):
    """流式调用大模型（自部署优先，MiMo 兜底，DeepSeek 二层兜底），逐段 yield 文本增量。

    流式模式不支持 response_format，输出格式由 system 提示词约束、
    由调用方解析。若当前 provider 已推送过内容则不再回退（回退会导致答案重复）。
    provider_sink 传入 list 时，成功产出内容的 provider 名称会被 append 进去，
    供调用方记录实际使用的模型标识。
    """
    messages = _llm_messages(prompt, system=system)
    providers = llm_config.available_providers()
    if not providers:
        raise RuntimeError("未配置任何大模型（SELF_LLM_API_KEY / DEEPSEEK_INTERNAL_BASE_URL / MIMO_API_KEY / DEEPSEEK_API_KEY）")
    last_error = None
    for provider in providers:
        try:
            produced = False
            for delta in _stream_provider(provider, messages, max_tokens):
                produced = True
                if provider_sink is not None and not provider_sink:
                    provider_sink.append(provider["name"])
                yield delta
            return
        except Exception as exc:
            if produced:
                raise
            last_error = exc
            current_app.logger.warning("%s 流式调用失败，尝试下一个大模型：%s", provider["name"], exc)
    raise RuntimeError(f"大模型流式调用失败：{last_error}")


def _provider_payload(provider, messages, max_tokens, json_mode, stream=False):
    """按 provider 能力构造 chat/completions 请求体。"""
    payload = {"model": provider["model"], "messages": messages, "temperature": 0.3}
    if json_mode and provider.get("json_mode"):
        payload["response_format"] = {"type": "json_object"}
    if max_tokens:
        payload["max_tokens"] = max_tokens
    if provider.get("disable_thinking"):
        # MiMo 关闭深度思考，直接输出答案，大幅降低首响应延迟
        payload["thinking"] = {"type": "disabled"}
    if provider.get("chat_template_kwargs"):
        # vLLM 部署的自部署 GLM：通过聊天模板参数关闭思考，
        # 响应从约 7s 降至 1s，且不再输出 reasoning_content
        payload["chat_template_kwargs"] = provider["chat_template_kwargs"]
    if stream:
        payload["stream"] = True
    return payload


def _complete_provider(provider, messages, max_tokens, json_mode):
    """非流式调用单个 provider：读取完整响应并返回 content 字符串。"""
    payload = _provider_payload(provider, messages, max_tokens, json_mode)
    try:
        with _open_llm_raw(provider["base_url"], payload, provider["api_key"]) as resp:
            data = json.loads(resp.read().decode("utf-8"))
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="replace")[:240]
        raise RuntimeError(f"知识服务返回 HTTP {exc.code}{(': ' + detail) if detail else ''}") from exc
    except urllib.error.URLError as exc:
        reason = getattr(exc, "reason", exc)
        if isinstance(reason, TimeoutError) or "timed out" in str(reason).lower():
            raise RuntimeError("知识服务响应超时，请稍后重试") from exc
        raise RuntimeError("无法连接知识服务，请检查网络连接或大模型 API 配置") from exc
    return data["choices"][0]["message"]["content"]


def _stream_provider(provider, messages, max_tokens):
    """流式调用单个 provider，逐段 yield 模型输出的文本增量。

    推理模型（MiMo、自部署 GLM 未关思考时）流中会先推送 reasoning_content
    增量（此处跳过），且可能出现空 choices 帧，需要防护。
    """
    payload = _provider_payload(provider, messages, max_tokens, json_mode=False, stream=True)
    try:
        with _open_llm_raw(provider["base_url"], payload, provider["api_key"]) as resp:
            for raw_line in resp:
                line = raw_line.decode("utf-8", errors="replace").strip()
                if not line.startswith("data:"):
                    continue
                data = line[len("data:"):].strip()
                if data == "[DONE]":
                    break
                try:
                    chunk = json.loads(data)
                except json.JSONDecodeError:
                    continue
                choices = chunk.get("choices") or []
                if not choices:
                    continue
                delta = (choices[0].get("delta") or {}).get("content")
                if delta:
                    yield delta
    except urllib.error.HTTPError as exc:
        detail = exc.read().decode("utf-8", errors="replace")[:240]
        raise RuntimeError(f"知识服务返回 HTTP {exc.code}{(': ' + detail) if detail else ''}") from exc
    except urllib.error.URLError as exc:
        reason = getattr(exc, "reason", exc)
        if isinstance(reason, TimeoutError) or "timed out" in str(reason).lower():
            raise RuntimeError("知识服务响应超时，请稍后重试") from exc
        raise RuntimeError("无法连接知识服务，请检查网络连接或大模型 API 配置") from exc
    except TimeoutError as exc:
        # 读取响应体阶段的超时不会包成 URLError，需要单独翻译。
        raise RuntimeError("知识服务响应超时，请稍后重试") from exc


def _llm_messages(prompt, system=None):
    messages = []
    if system:
        messages.append({"role": "system", "content": system})
    messages.append({"role": "user", "content": prompt})
    return messages


def _open_llm_raw(base_url, payload, api_key):
    """构造并发送 chat/completions 请求，返回未读取的响应对象。"""
    req = urllib.request.Request(
        base_url.rstrip("/") + "/chat/completions",
        data=json.dumps(payload).encode("utf-8"),
        headers={"Content-Type": "application/json",
                 "Authorization": f"Bearer {api_key}"},
    )
    # 系统环境可能配置了失效的本地代理（例如 127.0.0.1:9）。
    # 大模型请求显式直连，避免知识搜索因代理拒绝连接而失败。
    opener = urllib.request.build_opener(urllib.request.ProxyHandler({}))
    return opener.open(req, timeout=LLM_TIMEOUT)


def _ai_complete_tags_summary(file_path, original_name, external=False):
    """从文档中抽取文本，调用 DeepSeek 生成关键词和摘要。
    返回 {"tags": [...], "summary": "..."}，失败时抛异常。
    """
    text = _extract_text(file_path)
    context = text if text.strip() else f"文件名：{original_name}"
    prompt = (
        "你是金融研究报告助手。请根据下面的研究报告内容，提取关键词和摘要。"
        + ("同时识别报告原作者和发布机构；无法确认时返回空字符串。" if external else "")
        + "\n\n"
        "要求：\n"
        "1. 只输出一个JSON对象，键为 tags、summary、author、institution\n"
        "2. tags：3-6 个关键词，反映报告核心主题，用词简洁\n"
        "3. summary：不超过150字的中文摘要，概括报告的研究主题、主要内容或核心观点\n"
        "4. author 和 institution 仅填写文档明确出现的原作者/发布机构，不要把上传人当作作者\n"
        "5. 严格基于文档内容，不要编造\n\n"
        f"报告内容：\n{context}"
    )
    raw = _call_llm(prompt)
    data = json.loads(raw)
    tags = data.get("tags", [])
    if not isinstance(tags, list):
        tags = [str(tags)]
    tags = [str(t).strip() for t in tags if str(t).strip()][:8]
    summary = str(data.get("summary", "")).strip()
    return {
        "tags": tags,
        "summary": summary,
        "author": str(data.get("author", "")).strip()[:100],
        "institution": str(data.get("institution", "")).strip()[:120],
    }


def _report_meta_block(report):
    """报告元数据文本块，摘要与单篇问答共用。"""
    lines = [
        f"标题：{report.get('title') or ''}",
        f"作者：{report.get('sourceAuthor') or report.get('author') or ''}",
        f"机构：{report.get('sourceInstitution') or report.get('org') or ''}",
        f"报告种类：{REPORT_TYPE_LABELS.get(report.get('reportType') or '', '报告')}",
        f"研究主题：{REPORT_THEME_LABELS.get(report.get('theme') or '', '未分类')}",
        f"报告日期：{report.get('reportDate') or ''}",
    ]
    if report.get("summary"):
        lines.append(f"人工简介：{report['summary']}")
    return "\n".join(lines)


def _ai_summary_messages(report, text, style):
    """构造单篇报告摘要的 (system, user) 提示词。"""
    prompt = (
        f"【报告信息】\n{_report_meta_block(report)}\n\n"
        f"【报告全文】\n{text[:AI_SUMMARY_TEXT_CHARS]}\n\n"
        f"{AI_SUMMARY_STYLE_PROMPTS[style]}\n{AI_SUMMARY_COMMON_RULES}"
    )
    return AI_SUMMARY_SYSTEM, prompt


def _report_ask_messages(report, question):
    """构造单篇报告问答的 (system, user) 提示词：已缓存的 AI 摘要 + 全文。"""
    summaries = []
    for style in AI_SUMMARY_STYLES:
        row = _valid_summary_cache(report, style)
        if row:
            summaries.append(f"【AI{AI_SUMMARY_STYLES[style]['label']}】\n{row['content']}")
    file_path = report_file_path(report)
    text = _extract_text(file_path, max_chars=AI_SUMMARY_TEXT_CHARS) if file_path else ""
    system = (
        "你是研究报告阅读助手。仅基于给定报告的内容回答问题，不编造数据或结论；"
        "报告未涉及时明确说明“报告未提及”。回答简洁，用中文，可分点。"
    )
    prompt = (
        f"【报告信息】\n{_report_meta_block(report)}\n\n"
        f"【报告 AI 摘要】\n{''.join(summaries) if summaries else '（暂无）'}\n\n"
        f"【报告全文】\n{text.strip()[:AI_SUMMARY_TEXT_CHARS] if text.strip() else '（无法提取正文文本）'}\n\n"
        f"【问题】\n{question}"
    )
    return system, prompt


def _valid_summary_cache(report, style):
    """读取该报告该版本的摘要缓存；文件指纹（fileSha256）变化视为失效。"""
    row = store.get_report_summary(report.get("id"), style)
    if not row:
        return None
    if str(report.get("fileSha256") or "") != str(row.get("file_sha256") or ""):
        return None
    return row


def _summary_payload(row, style, user_id):
    generator = store.get_user(user_id) if user_id else None
    return {
        "style": style,
        "summary": row["content"],
        "generatedAt": row["updated_at"],
        "generatedByName": (generator or {}).get("name", ""),
        "model": row.get("model") or "",
    }


KNOWLEDGE_STOPWORDS = {
    "请问", "帮我", "告诉我", "什么", "如何", "哪些", "是否", "有没有", "最近", "近期",
    "目前", "当前", "相关", "报告", "研究", "分析", "观点", "情况", "主要", "内容",
    "影响", "原因", "变化", "方面", "关于", "以及", "可以", "一下", "今年", "最新",
    "主题", "完全", "存在", "不存在", "现在", "部门", "结论", "有什么", "报告",
    "报告么", "有么", "是否有", "里面", "这里", "其中",
}


def _knowledge_query_terms(question, vocabulary=""):
    """提取适合中文报告召回的主题词，降低“报告/研究/如何”等泛词的干扰。"""
    raw_terms = re.split(r"[\s,，。；;、？！?：:（）()]+", question)
    terms = {
        term.lower().strip() for term in raw_terms
        if len(term.strip()) >= 2 and term.strip() not in KNOWLEDGE_STOPWORDS
    }
    compact_query = re.sub(r"\s+", "", question).lower()
    subterms = set()
    for size in (2, 3, 4):
        for i in range(max(len(compact_query) - size + 1, 0)):
            piece = compact_query[i:i + size]
            if piece not in KNOWLEDGE_STOPWORDS and (not vocabulary or piece in vocabulary):
                subterms.add(piece)
    return compact_query, terms, subterms


# 知识搜索范围筛选：时间范围对应的天数，all/custom 由前端单独处理
KNOWLEDGE_PERIOD_DAYS = {"1m": 30, "3m": 90}
KNOWLEDGE_REPORT_TYPES = ("internal", "external", "research_visit", "roadshow")
KNOWLEDGE_FILTER_MAX_VALUES = 30  # 多选每个维度最多接受的选项数，防止恶意超长列表


def _parse_knowledge_filters(data):
    """解析知识搜索的范围筛选参数。

    时间范围为单选（1m/3m/all/custom + 自定义起止日期）；
    来源/种类/主题/人员支持多选，报告属性命中任一即可，
    空列表表示该维度不过滤。整体缺省为“过去一个月的全部报告”。
    """
    data = data or {}

    def selection(key, allowed, default):
        raw = data.get(key)
        if isinstance(raw, str):
            raw = [raw]
        if raw is None:
            return list(default)
        if not isinstance(raw, list):
            return list(default)
        return [str(item).strip() for item in raw
                if str(item).strip() in allowed][:KNOWLEDGE_FILTER_MAX_VALUES]

    period = str(data.get("period") or "").strip()
    if period not in ("1m", "3m", "all", "custom"):
        period = "1m"

    def parse_date(key):
        value = str(data.get(key) or "").strip()
        try:
            return datetime.strptime(value, "%Y-%m-%d").strftime("%Y-%m-%d")
        except ValueError:
            return ""

    date_from = date_to = ""
    if period == "custom":
        # 自定义区间为闭区间 [dateFrom, dateTo]，缺省的边界不限制。
        date_from = parse_date("dateFrom")
        date_to = parse_date("dateTo")
    elif period != "all":
        date_from = (datetime.now(CST) - timedelta(days=KNOWLEDGE_PERIOD_DAYS[period])).strftime("%Y-%m-%d")

    raw_types = data.get("reportTypes")
    if isinstance(raw_types, str):
        raw_types = [raw_types]
    if raw_types is None:
        report_types = []
    elif isinstance(raw_types, list):
        report_types = [str(item).strip() for item in raw_types
                        if str(item).strip() in KNOWLEDGE_REPORT_TYPES][:KNOWLEDGE_FILTER_MAX_VALUES]
    else:
        report_types = []

    raw_authors = data.get("authors")
    if isinstance(raw_authors, str):
        raw_authors = [raw_authors]
    authors = ([str(item).strip()[:50] for item in raw_authors if str(item).strip()]
               if isinstance(raw_authors, list) else [])[:KNOWLEDGE_FILTER_MAX_VALUES]

    return {
        "period": period, "date_from": date_from, "date_to": date_to,
        "report_types": report_types,
        "categories": selection("categories", REPORT_CATEGORY_LABELS, ()),
        "themes": selection("themes", REPORT_THEMES, ()),
        "authors": authors,
    }


def _report_matches_knowledge_filters(report, filters):
    """判断报告是否落在知识搜索筛选范围内；filters 为空表示不过滤。

    来源/种类/主题/人员为多选：报告属性命中已选列表中的任一值即通过，
    列表为空表示该维度不过滤。
    """
    if not filters:
        return True
    report_types = filters.get("report_types") or []
    if report_types and str(report.get("reportType", "internal")) not in report_types:
        return False
    categories = filters.get("categories") or []
    if categories and str(report.get("category", "")) not in categories:
        return False
    themes = filters.get("themes") or []
    if themes and str(report.get("theme", "")) not in themes:
        return False
    authors = filters.get("authors") or []
    if authors:
        # 外部报告的署名作者在 sourceAuthor，内部报告在 author。
        names = {str(report.get("author", "")), str(report.get("sourceAuthor", ""))}
        if not names.intersection(authors):
            return False
    if filters.get("date_from") or filters.get("date_to"):
        date_str = str(report.get("reportDate") or report.get("uploadedAt", ""))[:10]
        if not re.fullmatch(r"\d{4}-\d{2}-\d{2}", date_str):
            return False
        if filters["date_from"] and date_str < filters["date_from"]:
            return False
        if filters["date_to"] and date_str > filters["date_to"]:
            return False
    return True


def _knowledge_no_match_answer(filters=None):
    """筛选范围内无候选时的回答文案，提示用户放宽条件而不是报告库为空。"""
    if filters:
        return "当前筛选范围内未找到相关报告，可尝试扩大时间范围或放宽筛选条件。"
    return "未找到相关报告"


KNOWLEDGE_SYNONYM_GROUPS = (
    ("城投", "平台公司", "地方政府融资平台"),
    ("地产", "房地产", "房企"),
    ("资金面", "流动性", "银行间资金"),
    ("债市", "债券市场", "固定收益市场"),
    ("利差", "信用利差", "spread"),
    ("收益率", "到期收益率", "yield"),
    ("二永债", "二级资本债", "永续债"),
)


def _knowledge_metadata(report):
    """Build a stable searchable metadata block shared by indexing and ranking."""
    theme_key = str(report.get("theme", ""))
    category_key = str(report.get("category", ""))
    return "\n".join(filter(None, [
        str(report.get("title", "")), str(report.get("summary", "")),
        str(report.get("recommendation", "")), " ".join(report.get("tags") or []),
        str(report.get("author", "")), str(report.get("sourceAuthor", "")),
        str(report.get("sourceInstitution", "")), str(report.get("org", "")),
        theme_key, REPORT_THEME_LABELS.get(theme_key, ""),
        category_key, REPORT_CATEGORY_LABELS.get(category_key, ""),
    ]))


def _knowledge_fingerprint(report):
    """Fingerprint searchable metadata plus immutable upload file state."""
    path = report_file_path(report)
    file_state = None
    if path:
        stat = os.stat(path)
        file_state = [stat.st_size, stat.st_mtime_ns]
    payload = {
        "version": KNOWLEDGE_VECTOR_VERSION,
        "metadata": _knowledge_metadata(report),
        "file": file_state,
    }
    raw = json.dumps(payload, ensure_ascii=False, sort_keys=True).encode("utf-8")
    return hashlib.sha256(raw).hexdigest()


def _knowledge_vector_tokens(value):
    """Tokenize Chinese text into character n-grams and ASCII terms."""
    text = unicodedata.normalize("NFKC", str(value or "")).lower()
    tokens = re.findall(r"[a-z][a-z0-9_+.-]{1,}|\d+(?:\.\d+)?", text)
    for run in re.findall(r"[\u3400-\u9fff]+", text):
        for size in (2, 3, 4):
            tokens.extend(run[index:index + size] for index in range(len(run) - size + 1))
    for group in KNOWLEDGE_SYNONYM_GROUPS:
        if any(term in text for term in group):
            tokens.extend(group)
    return tokens


def _knowledge_vector(value):
    """Create a deterministic, normalized feature-hashed vector without network calls."""
    counts = Counter(_knowledge_vector_tokens(value))
    vector = [0.0] * KNOWLEDGE_VECTOR_DIM
    for token, count in counts.items():
        digest = hashlib.blake2b(token.encode("utf-8"), digest_size=8,
                                 person=b"ikb-vector-v1").digest()
        hashed = int.from_bytes(digest, "little")
        index = hashed % KNOWLEDGE_VECTOR_DIM
        sign = -1.0 if hashed & (1 << 63) else 1.0
        length_weight = 1.0 + min(max(len(token) - 2, 0), 3) * 0.12
        vector[index] += sign * (1.0 + math.log(count)) * length_weight
    norm = math.sqrt(sum(item * item for item in vector))
    if norm:
        vector = [item / norm for item in vector]
    return vector


def _knowledge_pack_vector(vector):
    return struct.pack(f"<{KNOWLEDGE_VECTOR_DIM}f", *vector)


def _knowledge_unpack_vector(payload):
    expected = KNOWLEDGE_VECTOR_DIM * 4
    if not payload or len(payload) != expected:
        return None
    return struct.unpack(f"<{KNOWLEDGE_VECTOR_DIM}f", payload)


def _knowledge_text_chunks(text):
    value = re.sub(r"\n{3,}", "\n\n", str(text or "")).strip()
    if not value:
        return []
    chunks = []
    start = 0
    while start < len(value):
        end = min(start + KNOWLEDGE_CHUNK_CHARS, len(value))
        if end < len(value):
            boundary = max(value.rfind("\n", start + KNOWLEDGE_CHUNK_CHARS // 2, end),
                           value.rfind("。", start + KNOWLEDGE_CHUNK_CHARS // 2, end))
            if boundary > start:
                end = boundary + 1
        chunks.append(value[start:end].strip())
        if end >= len(value):
            break
        start = max(end - KNOWLEDGE_CHUNK_OVERLAP, start + 1)
    return [item for item in chunks if item]


def _index_knowledge_report(report, force=False):
    """Extract, chunk and persist one report; return True when rebuilt."""
    report_id = report["id"]
    fingerprint = _knowledge_fingerprint(report)
    current = store.knowledge_index_fingerprints([report_id]).get(report_id)
    if (not force and current and current.get("fingerprint") == fingerprint and
            current.get("vector_version") == KNOWLEDGE_VECTOR_VERSION):
        return False
    with _knowledge_index_lock:
        current = store.knowledge_index_fingerprints([report_id]).get(report_id)
        if (not force and current and current.get("fingerprint") == fingerprint and
                current.get("vector_version") == KNOWLEDGE_VECTOR_VERSION):
            return False
        metadata = _knowledge_metadata(report)
        path = report_file_path(report)
        body = _extract_text(path, max_chars=KNOWLEDGE_INDEX_TEXT_CHARS) if path else ""
        chunks = _knowledge_text_chunks(body) or [metadata or report.get("title", "未命名报告")]
        store.replace_knowledge_chunks(
            report_id, fingerprint, KNOWLEDGE_VECTOR_VERSION,
            [{"content": chunk, "vector": _knowledge_pack_vector(
                _knowledge_vector(metadata + "\n" + chunk))} for chunk in chunks],
        )
    return True


def rebuild_knowledge_vector_index(force=False, progress=None):
    """Operational backfill entry point used by the isolated test and deployment script."""
    reports = store.reports()
    rebuilt = 0
    for index, report in enumerate(reports, 1):
        rebuilt += int(_index_knowledge_report(report, force=force))
        if progress:
            progress(index, len(reports), report)
    return {**store.knowledge_index_stats(), "rebuilt": rebuilt, "total": len(reports),
            "vectorVersion": KNOWLEDGE_VECTOR_VERSION}


def _knowledge_lexical_score(query, terms, subterms, report, content):
    title = str(report.get("title", ""))
    metadata_compact = re.sub(r"\s+", "", _knowledge_metadata(report)).lower()
    title_compact = re.sub(r"\s+", "", title).lower()
    content_compact = re.sub(r"\s+", "", content).lower()
    full_compact = metadata_compact + content_compact
    score = 0
    if query and len(query) >= 3 and query in full_compact:
        score += 18
    for term in terms:
        if term in title_compact:
            score += 14
        elif term in metadata_compact:
            score += 6
        if term in content_compact:
            score += 4
    for sub in subterms:
        if sub in metadata_compact:
            score += max(6, min(len(sub) * 2, 8))
        elif sub in content_compact:
            score += max(3, min(len(sub), 4))
    return score


def _iter_knowledge_candidates(question, limit=6, filters=None):
    """Hybrid vector/keyword retrieval with persistent chunk embeddings.

    Each report is extracted only when its source fingerprint changes. Subsequent searches
    rank persisted vectors and therefore avoid reparsing every PDF/Office file.
    filters：可选的范围筛选（_parse_knowledge_filters 的结果），检索前先过滤报告。
    """
    reports = [report for report in store.reports()
               if _report_matches_knowledge_filters(report, filters)]
    metadata_corpus = "".join(
        re.sub(r"\s+", "", " ".join([
            str(report.get("title", "")), str(report.get("summary", "")),
            str(report.get("recommendation", "")), " ".join(report.get("tags") or []),
            str(report.get("author", "")), str(report.get("theme", "")),
            REPORT_THEME_LABELS.get(str(report.get("theme", "")), ""),
            str(report.get("category", "")), REPORT_CATEGORY_LABELS.get(str(report.get("category", "")), ""),
        ])).lower()
        for report in reports
    )
    query, terms, subterms = _knowledge_query_terms(question, metadata_corpus)
    reports_by_id = {report["id"]: report for report in reports}
    indexed = store.knowledge_index_fingerprints(reports_by_id)
    for idx, report in enumerate(reports, 1):
        current = indexed.get(report["id"])
        fingerprint = _knowledge_fingerprint(report)
        if (not current or current.get("fingerprint") != fingerprint or
                current.get("vector_version") != KNOWLEDGE_VECTOR_VERSION):
            _index_knowledge_report(report)
        yield idx, len(reports)
    if not reports:
        return []

    query_vector = _knowledge_vector(question)
    ranked_by_report = {}
    title_affinity = {}
    for report_id, report in reports_by_id.items():
        title = re.sub(r"\s+", "", str(report.get("title", ""))).lower()
        title_vector = _knowledge_vector(title)
        title_cosine = sum(left * right for left, right in zip(query_vector, title_vector))
        title_hits = sum((len(term) - 1) ** 2 for term in subterms if term in title)
        title_affinity[report_id] = max(title_cosine, 0) * 55.0 + min(title_hits * 1.8, 45.0)
    for row in store.knowledge_chunks(reports_by_id):
        vector = _knowledge_unpack_vector(row["vector"])
        if vector is None:
            continue
        report = reports_by_id.get(row["report_id"])
        if not report:
            continue
        content = row["content"]
        cosine = sum(left * right for left, right in zip(query_vector, vector))
        lexical = _knowledge_lexical_score(query, terms, subterms, report, content)
        # Hybrid score keeps exact title/metadata matches decisive while vector similarity
        # recalls wording variants and ranks the best source passage.
        hybrid = (cosine * 100.0 + min(lexical, 40) * 1.5 +
                  title_affinity.get(report["id"], 0.0))
        if lexical < 6 and cosine < 0.115:
            continue
        entry = ranked_by_report.setdefault(report["id"], {
            "score": float("-inf"), "vector_score": cosine, "lexical_score": lexical,
            "published_at": str(report.get("reportDate") or report.get("uploadedAt", "")),
            "report": report, "chunks": [],
        })
        entry["score"] = max(entry["score"], hybrid)
        entry["vector_score"] = max(entry["vector_score"], cosine)
        entry["lexical_score"] = max(entry["lexical_score"], lexical)
        entry["chunks"].append((hybrid, content))

    ranked = sorted(ranked_by_report.values(),
                    key=lambda item: (item["score"], item["published_at"]), reverse=True)
    # 同标题报告只保留一份，优先保留相关度更高、发布时间更晚的版本。
    unique = []
    seen_titles = set()
    for item in ranked:
        title_key = re.sub(r"\s+", "", str(item["report"].get("title", "未命名报告"))).lower()
        if title_key in seen_titles:
            continue
        seen_titles.add(title_key)
        unique.append(item)
        if len(unique) >= limit:
            break

    candidates = []
    for item in unique:
        report = item["report"]
        published_at = item["published_at"]
        passages = []
        for _, passage in sorted(item["chunks"], reverse=True):
            if passage in passages:
                continue
            passages.append(passage)
            if sum(len(part) for part in passages) >= 5200 or len(passages) >= 3:
                break
        candidates.append({
            "id": report["id"],
            "title": report.get("title", "未命名报告"),
            "author": report.get("author", ""),
            "published_at": published_at[:10] if published_at else "原文未提及",
            "theme": REPORT_THEME_LABELS.get(report.get("theme"), report.get("theme", "原文未提及")),
            "category": REPORT_CATEGORY_LABELS.get(report.get("category"), report.get("category", "原文未提及")),
            "content": "\n\n".join(passages)[:5500],
        })
    return candidates


def _knowledge_candidates(question, limit=6, filters=None):
    """检索正文与元信息，过滤弱相关项，并返回带完整元数据的上下文候选。"""
    retrieval = _iter_knowledge_candidates(question, limit=limit, filters=filters)
    while True:
        try:
            next(retrieval)
        except StopIteration as stop:
            return stop.value or []


KNOWLEDGE_SOURCE_MARKER = "===SOURCE_IDS==="

KNOWLEDGE_INTENT_RETRIEVAL = "report_retrieval"
KNOWLEDGE_INTENT_GENERAL = "general_work"

_KNOWLEDGE_GENERAL_INTENT_PATTERNS = (
    # 数量、占比、分布等统计任务，需要跨报告汇总而不是逐篇摘要。
    r"统计|计数|多少篇|数量|占比|比例|分布|频次|排名|排行|均值|平均|中位数",
    # 多报告综合与观点关系分析。
    r"比较|对比|异同|共同点|共识|分歧|一致性|交叉验证|观点演变|趋势演变|横向|纵向",
    r"汇总|归纳|分类|聚类|矩阵|整体观点|综合观点|全部报告|所有报告|多份报告|各份报告|不同报告",
    # 基于材料继续完成工作成果或发散任务。
    r"撰写|起草|拟一份|写一份|生成.*(?:提纲|框架|清单|表格|方案)|制作.*(?:提纲|框架|清单|表格)",
    r"头脑风暴|发散|启示|研究方向|下一步|情景推演|策略框架|行动建议|工作建议",
)


def _knowledge_intent(question):
    """本地判断知识问答意图，避免为路由提示词额外消耗一次模型调用。

    统计、跨报告综合、成果撰写和发散任务进入通用框架；其余事实型问题、
    报告查找和核心观点查询继续使用原有报告检索框架。无法明确判断时回退
    到原有框架，以保持既有问答行为稳定。
    """
    text = re.sub(r"\s+", "", str(question or "")).lower()
    # “归纳这篇报告的核心观点”仍属于原有核心观点查询；只有同时出现统计、
    # 跨报告比较或成果制作等目标时，才进入通用工作框架。
    if re.search(r"核心观点|主要观点|核心结论|主要结论", text):
        advanced_core_task = re.search(
            r"统计|计数|多少篇|数量|占比|分布|频次|排名|比较|对比|异同|共同点|共识|分歧|"
            r"一致性|交叉验证|演变|全部报告|所有报告|多份报告|各份报告|不同报告|矩阵|聚类|"
            r"撰写|起草|拟一份|写一份|生成|制作|头脑风暴|发散|推演|建议",
            text,
        )
        if not advanced_core_task:
            return KNOWLEDGE_INTENT_RETRIEVAL
    if any(re.search(pattern, text) for pattern in _KNOWLEDGE_GENERAL_INTENT_PATTERNS):
        return KNOWLEDGE_INTENT_GENERAL
    return KNOWLEDGE_INTENT_RETRIEVAL


def _knowledge_qa_messages(question, candidates, intent=None):
    """构建知识问答的 system 与用户 prompt，供流式与非流式两条路径共用。

    输出格式为“答案正文 + 单独一行的标记 + JSON 数组”，使流式接口可以
    把标记之前的正文实时推送给前端，标记之后只留机器可读的引用 ID。
    """
    intent = intent or _knowledge_intent(question)
    ordered_candidates = sorted(candidates, key=lambda item: item["published_at"], reverse=True)
    content_limit = 3200 if intent == KNOWLEDGE_INTENT_GENERAL else None
    material = "\n\n".join(
        f"[REPORT_ID:{item['id']}]\n"
        f"报告完整标题：{item['title']}\n"
        f"报告时间：{item['published_at']}\n"
        f"报告元数据方向：{item['theme']} / {item['category']}\n"
        f"报告作者：{item['author'] or '原文未提及'}\n"
        f"报告正文：\n{item['content'][:content_limit] if content_limit else item['content']}"
        for item in ordered_candidates
    )
    if intent == KNOWLEDGE_INTENT_GENERAL:
        system = (
            "你是一个专业的内部研究知识工作助手。你的任务不是机械罗列报告，而是仅依据给定的"
            "【参考上下文】，完成用户要求的统计、归纳、比较、组织材料、形成框架或其他研究工作。\n\n"
            "严格规则：\n"
            "1. 仅使用【参考上下文】中的事实和观点，禁止引入外部资料或编造；可以进行必要的归纳推理，"
            "但必须明确区分“报告原文观点”和“基于多份报告的综合判断”。\n"
            "2. 先识别用户真正要完成的工作，再选择最合适的结构；可使用小标题、项目符号、编号或表格，"
            "不强制套用逐篇报告摘要模板。\n"
            "3. 统计任务必须先说明统计口径、样本范围和样本数量，所有数字都须能由上下文逐项核验；"
            "本次上下文是相关报告样本时，不得声称统计覆盖整个报告库。\n"
            "4. 比较或归纳任务应优先呈现共识、分歧、变化和证据，并注明支撑结论的报告标题或时间。\n"
            "5. 提纲、框架、建议或发散任务可以在报告观点之上继续组织，但不得把延伸建议伪装成报告原文结论。\n"
            "6. 合并重复信息，保持专业、简洁、可直接用于工作；上下文不足以完成任务时，明确指出缺口。\n"
            "7. 只引用实际支持答案的报告 ID，不要引用仅关键词相似但没有提供证据的报告。\n\n"
            "输出格式：先直接输出适合该任务的中文答案正文，不要添加无意义的开场白或结束语；"
            f"然后另起一行只输出 {KNOWLEDGE_SOURCE_MARKER}，"
            "紧接着再起一行输出一个 JSON 数组，数组内只包含实际支持答案的报告 ID，"
            "例如 [\"r001\",\"r002\"]；没有引用任何报告时输出 []。除此之外不要输出任何内容。"
        )
        task_hint = (
            f"系统已将该问题识别为综合工作任务。本次提供 {len(ordered_candidates)} 篇相关报告样本，"
            "请根据用户目标灵活组织答案。"
        )
    else:
        system = (
            "你是一个专业的知识库报告检索与分析助手。你的核心任务是根据用户问题，"
            "从给定的【参考上下文】中精准筛选相关报告，并按指定结构输出摘要。\n\n"
            "严格规则：\n"
            "1. 仅使用【参考上下文】回答，禁止使用外部知识、常识补全或编造；上下文没有实质相关报告时，answer 必须严格为“未找到相关报告”。\n"
            "2. 客观中立，忠实还原报告原文观点，不添加个人评价、推测或总结性升华。\n"
            "3. 每篇报告必须包含发布时间、主要方向、主要观点；上下文缺失时对应字段写“原文未提及”。\n"
            "4. 合并重复报告，按发布时间倒序排列；只引用实际入选且支持答案的报告 ID。\n"
            "5. 逐一判断报告与问题是否实质相关，剔除仅关键词匹配但内容无关的报告。\n"
            "6. 每条主要观点精炼至 100 字以内，最多输出 3 条；不要输出上下文没有明确支持的数字、因果关系或结论。\n\n"
            "answer 必须严格使用以下模板，不要添加任何开场白、结束语或模板外内容：\n"
            "共检索到 {N} 篇相关报告：\n\n"
            "### 1. 《{报告完整标题}》\n"
            "- **发布时间**：{YYYY-MM-DD / YYYY年MM月}\n"
            "- **主要方向**：{用1句话概括报告研究的核心领域或主题}\n"
            "- **主要观点**：\n"
            "  - {核心观点1，精炼至100字以内}\n"
            "  - {核心观点2，精炼至100字以内}\n"
            "  - {核心观点3（如有），精炼至100字以内}\n\n"
            "依此类推。\n\n"
            "输出格式：先直接输出 answer 正文（严格按上述模板，不要用引号或代码块包裹），"
            f"然后另起一行只输出 {KNOWLEDGE_SOURCE_MARKER}，"
            "紧接着再起一行输出一个 JSON 数组，数组内只包含实际入选且支持答案的报告 ID，"
            "例如 [\"r001\",\"r002\"]；没有引用任何报告时输出 []。除此之外不要输出任何内容。"
        )
        task_hint = "系统已将该问题识别为报告检索或核心观点查询，请严格筛选相关报告。"
    prompt = f"{task_hint}\n\n用户问题：{question}\n\n【参考上下文】\n{material}"
    return system, prompt


def _parse_knowledge_answer(raw, candidates):
    """解析模型输出：正文 + ===SOURCE_IDS=== 标记 + JSON 数组。

    兼容模型偶尔仍按旧 JSON 对象格式（{\"answer\":..., \"source_ids\":[...]}）
    输出的情况，保证非流式路径的旧行为不受影响。
    """
    text = str(raw or "").strip()
    source_ids = []
    if KNOWLEDGE_SOURCE_MARKER in text:
        text, ids_raw = text.split(KNOWLEDGE_SOURCE_MARKER, 1)
        text = text.strip()
        match = re.search(r"\[.*?\]", ids_raw, re.S)
        if match:
            try:
                parsed = json.loads(match.group(0))
            except json.JSONDecodeError:
                parsed = None
            if isinstance(parsed, list):
                source_ids = parsed
    else:
        data = None
        try:
            data = json.loads(text)
        except json.JSONDecodeError:
            found = re.search(r"\{.*\}", text, re.S)
            if found:
                try:
                    data = json.loads(found.group(0))
                except json.JSONDecodeError:
                    data = None
        if isinstance(data, dict):
            ids = data.get("source_ids", [])
            source_ids = ids if isinstance(ids, list) else [ids]
            text = str(data.get("answer", "")).strip()
    source_map = {item["id"]: item for item in candidates}
    sources = [
        {"id": rid, "title": source_map[rid]["title"], "author": source_map[rid]["author"],
         "publishedAt": source_map[rid]["published_at"]}
        for rid in (str(item) for item in source_ids) if rid in source_map
    ]
    sources.sort(key=lambda item: item["publishedAt"], reverse=True)
    answer = text
    if not answer or answer == "报告库中暂无可用于回答的报告。":
        answer = "未找到相关报告"
    if answer == "未找到相关报告":
        sources = []
    return {"answer": answer, "sources": sources}


def _answer_knowledge_question(question, filters=None):
    intent = _knowledge_intent(question)
    candidate_limit = 12 if intent == KNOWLEDGE_INTENT_GENERAL else 6
    candidates = _knowledge_candidates(question, limit=candidate_limit, filters=filters)
    if not candidates:
        return {"answer": _knowledge_no_match_answer(filters), "sources": []}
    system, prompt = _knowledge_qa_messages(question, candidates, intent=intent)
    max_tokens = 3600 if intent == KNOWLEDGE_INTENT_GENERAL else 2800
    raw = _call_llm(prompt, system=system, max_tokens=max_tokens, json_mode=False)
    return _parse_knowledge_answer(raw, candidates)


def json_error(message, code=400):
    return jsonify({"error": message}), code


# --------------------------------------------------------------------------- #
# 静态文件 & 入口路由
# --------------------------------------------------------------------------- #
def html_page(filename):
    html = Path(BASE_DIR, filename).read_text(encoding="utf-8")
    html = html.replace("__CSRF_TOKEN__", csrf_token())
    return Response(html, content_type="text/html; charset=utf-8")


@app.route("/")
def index():
    return html_page("index.html")


@app.route("/admin")
def admin_page():
    return html_page("admin.html")


@app.route("/admin.html")
def admin_html_page():
    """兼容通过 HTTP 直接访问 admin.html 的入口。"""
    return html_page("admin.html")


@app.route("/index.html")
def index_html_page():
    """兼容通过 HTTP 直接访问 index.html 的入口。"""
    return html_page("index.html")


@app.route("/<path:path>")
def static_files(path):
    # 仅开放前端静态资源目录，避免数据库、上传文件和源码被直接下载。
    normalized = path.replace("\\", "/").lstrip("/")
    parts = normalized.split("/")
    if not parts or parts[0] not in {"assets", "css", "js"} or any(part in {"", ".", ".."} for part in parts):
        return ("Not Found", 404)
    full_path = os.path.join(BASE_DIR, path)
    if os.path.isfile(full_path):
        return send_from_directory(BASE_DIR, path)
    return ("Not Found", 404)


# --------------------------------------------------------------------------- #
# 普通用户认证 API
# --------------------------------------------------------------------------- #
@app.route("/api/login", methods=["POST"])
def api_login():
    data = request.get_json(silent=True) or {}
    username = data.get("username", "")
    password = data.get("password", "")
    user = store.find_user_by_login(username)
    if not user or not check_password_hash(user.get("password_hash", ""), password):
        return json_error("账号或密码不正确", 401)
    session["internal_knowledge_base_user_id"] = user["id"]
    return jsonify({"user": public_user(user)})


@app.route("/api/logout", methods=["POST"])
def api_logout():
    session.pop("internal_knowledge_base_user_id", None)
    return jsonify({"ok": True})


@app.route("/api/me", methods=["GET"])
def api_me():
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    return jsonify({"user": public_user(user)})


@app.route("/api/change-password", methods=["POST"])
def api_change_password():
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    data = request.get_json(silent=True) or {}
    old_pw = data.get("oldPassword", "")
    new_pw = data.get("newPassword", "")
    if not check_password_hash(user.get("password_hash", ""), old_pw):
        return json_error("当前密码不正确", 400)
    if len(new_pw) < 6:
        return json_error("新密码至少需要 6 位", 400)
    if new_pw == old_pw:
        return json_error("新密码不能与当前密码相同", 400)
    store.update_user(user["id"], {"password_hash": generate_password_hash(new_pw)})
    return jsonify({"ok": True})


# --------------------------------------------------------------------------- #
# 报告 API
# --------------------------------------------------------------------------- #
@app.route("/api/reports", methods=["GET"])
def api_reports_list():
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    schedule_by_id = {item["id"]: item for item in store.all_roadshow_items()}
    reports = []
    for report in store.reports():
        item = public_report(report)
        item.update(store.report_engagement(report["id"], user["id"]))
        # 路演报告附带关联路演安排摘要，供详情弹窗展示与手工匹配
        schedule_id = str(report.get("roadshowScheduleId") or "")
        if report.get("reportType") == "roadshow" and schedule_id and schedule_id in schedule_by_id:
            schedule = schedule_by_id[schedule_id]
            item["roadshowSchedule"] = {
                "id": schedule_id,
                "eventTime": schedule.get("event_time", ""),
                "endTime": schedule.get("end_time", ""),
                "presenter": schedule.get("presenter", ""),
                "topic": schedule.get("topic", ""),
                "institution": schedule.get("institution", ""),
            }
        reports.append(item)
    return jsonify({"reports": reports})


@app.route("/api/report-authors", methods=["GET"])
def api_report_authors():
    """供行政角色上传报告时选择署名作者。"""
    user = require_role("admin")
    if not user:
        return json_error("仅行政角色可选择报告作者", 403)
    return jsonify({"authors": [public_user(item) for item in store.users()]})


def _sha256_stream(stream) -> str:
    digest = hashlib.sha256()
    stream.seek(0)
    for chunk in iter(lambda: stream.read(1024 * 1024), b""):
        digest.update(chunk)
    stream.seek(0)
    return digest.hexdigest()


def _find_duplicate_report(sha256_hex: str):
    """按内容哈希查找未删除的既有报告；历史文件没有存哈希时读盘计算并回填。"""
    for report in store.reports():
        if not report.get("fileStored"):
            continue
        digest = str(report.get("fileSha256") or "")
        if not digest:
            stored_name = str(report.get("fileUrl") or "").split("/")[-1]
            path = os.path.join(UPLOAD_DIR, stored_name)
            if not stored_name or not os.path.isfile(path):
                continue
            try:
                with open(path, "rb") as fh:
                    digest = _sha256_stream(fh)
            except OSError:
                continue
            try:
                store.update_report(report["id"], {"fileSha256": digest})
            except Exception:
                current_app.logger.exception("回填报告文件哈希失败：%s", report.get("id"))
        if digest == sha256_hex:
            return report
    return None


@app.route("/api/reports", methods=["POST"], defaults={"report_scope": None})
@app.route("/api/reports/<report_scope>", methods=["POST"])
def api_reports_upload(report_scope=None):
    """单篇报告上传。FormData:
       - file: 单个文件
       - meta: JSON 字符串 {reportType, category, org, reportDate, summary,
               recommendation, tags, titles:{filename:title}, authorId}

       authorId 仅对行政角色生效；其他角色始终以当前登录用户作为作者。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)

    files = request.files.getlist("file") or request.files.getlist("files")
    files = [f for f in files if f and f.filename]
    if not files:
        return json_error("请选择要上传的报告文件", 400)
    if len(files) > 1:
        return json_error("一次只能上传一篇报告", 400)

    meta_raw = request.form.get("meta", "{}")
    try:
        meta = json.loads(meta_raw)
    except json.JSONDecodeError:
        return json_error("元信息格式错误", 400)

    category = meta.get("category", "").strip()
    theme = meta.get("theme", "").strip()
    report_type = meta.get("reportType", "internal").strip()
    if report_scope not in (None, "internal", "external", "research_visit", "roadshow"):
        return json_error("报告上传入口无效", 404)
    if report_scope and report_type != report_scope:
        return json_error("该入口只能上传对应类型的报告", 400)
    org = meta.get("org", "").strip()
    report_date = meta.get("reportDate", "").strip()
    summary = meta.get("summary", "").strip()
    recommendation = meta.get("recommendation", "").strip()[:300]
    source_author = meta.get("sourceAuthor", "").strip()[:100]
    source_institution = meta.get("sourceInstitution", "").strip()[:120]
    tags_raw = meta.get("tags", "")
    titles = meta.get("titles", {}) or {}

    # 打分部门：仅深度报告生效（月报打分已关闭），默认两部门均打分
    scoring_orgs_raw = meta.get("scoringOrgs", [])
    if not isinstance(scoring_orgs_raw, list):
        scoring_orgs_raw = [scoring_orgs_raw]
    scoring_orgs = [org for org in scoring_orgs_raw if org in SCORING_ORGS]
    if report_type == "internal" and category == "deep":
        if not scoring_orgs:
            return json_error("请至少选择一个打分部门", 400)
    else:
        scoring_orgs = list(SCORING_ORGS)

    # 路演安排表关联：从路演日程一键上传报告时携带，便于回溯关联；
    # 未携带时尝试自动匹配（规则 + 大模型，匹配不上留空，可后续手工匹配）
    roadshow_schedule_id = ""
    roadshow_match_method = ""
    if report_type == "roadshow":
        roadshow_schedule_id = str(meta.get("roadshowScheduleId", "")).strip()
        if roadshow_schedule_id and not store.get_roadshow_item(roadshow_schedule_id):
            return json_error("关联的路演安排不存在，请刷新后重试", 400)
        if roadshow_schedule_id:
            roadshow_match_method = "manual"

    author = user
    if user.get("role") == "admin":
        author_id = str(meta.get("authorId", "")).strip()
        if not author_id:
            return json_error("请选择报告作者", 400)
        author = store.get_user(author_id)
        if not author:
            return json_error("请选择有效的报告作者", 400)

    if report_type not in ("internal", "external", "research_visit", "roadshow"):
        return json_error("请选择有效的报告类型", 400)
    if report_type == "external":
        # 外部报告不使用内部报告分类字段，内部以 other 兼容旧数据结构保存。
        category = "other"
    elif report_type == "internal" and category not in ("weekly", "monthly", "deep", "other"):
        return json_error("请选择有效的报告分类", 400)
    if theme not in REPORT_THEMES:
        return json_error("请选择有效的报告主题", 400)
    if org not in ("资产配置部", "固收中心"):
        return json_error("请选择有效的所属部门", 400)
    if not report_date:
        return json_error("请选择报告日期", 400)
    if report_type in ("external", "roadshow") and not source_author:
        return json_error("请填写报告作者", 400)
    if report_type in ("external", "roadshow") and not source_institution:
        return json_error("请填写报告机构", 400)

    tags = [t.strip() for t in re.split(r"[,，]", tags_raw) if t.strip()][:8]

    created = []
    now_iso = _now_iso()
    for f in files:
        original = f.filename
        ext = os.path.splitext(original)[1].lower()
        if ext not in ALLOWED_UPLOAD_EXTS:
            continue
        # 读取文件内容检查大小
        f.stream.seek(0, 2)
        size = f.stream.tell()
        f.stream.seek(0)
        if size > MAX_UPLOAD_SIZE:
            return json_error(f"文件 {original} 超过 100MB 限制", 400)

        # 内容级查重：按文件字节计算 SHA-256 与既有报告比对（不只看文件名），
        # 历史报告首次比对时懒回填哈希；完全相同的内容直接拦截本次上传。
        sha256_hex = _sha256_stream(f.stream)
        duplicate = _find_duplicate_report(sha256_hex)
        if duplicate is not None:
            existing_date = str(duplicate.get("uploadedAt") or "")[:10]
            existing_uploader = duplicate.get("uploadedByName") or duplicate.get("author") or "未知"
            return json_error(
                f"该文件内容与已上传的《{duplicate.get('title') or duplicate.get('fileName') or '未命名报告'}》"
                f"（{existing_date} 由 {existing_uploader} 上传）完全相同，已拦截重复上传",
                400,
            )

        report_id = f"report-upload-{int(time.time() * 1000)}-{uuid.uuid4().hex[:8]}"
        save_name = stored_filename(report_id, original)
        save_path = os.path.join(UPLOAD_DIR, save_name)
        f.save(save_path)

        title = (titles.get(original) or "").strip() or os.path.splitext(original)[0]
        # 一键上传未携带关联时，按报告信息自动匹配路演安排（失败留空，不阻断上传）
        if report_type == "roadshow" and not roadshow_schedule_id:
            try:
                roadshow_schedule_id, roadshow_match_method = _roadshow_auto_match(
                    report_date, source_author, source_institution, title)
            except Exception:
                current_app.logger.exception("路演报告自动匹配失败：%s", report_id)
        report = {
            "id": report_id,
            "title": title,
            "author": author["name"],
            "authorId": author["id"],
            # 上传人始终记录实际操作者（行政代传时上传人=行政，署名作者不变）
            "uploadedById": user["id"],
            "uploadedByName": user["name"],
            "sourceAuthor": source_author if report_type in ("external", "roadshow") else "",
            "sourceInstitution": source_institution if report_type in ("external", "roadshow") else "",
            "org": org,
            "category": category,
            "theme": theme,
            "reportType": report_type,
            "reportDate": report_date,
            "uploadedAt": now_iso,
            "summary": summary,
            "recommendation": recommendation if report_type == "external" else "",
            "tags": tags,
            "fileName": original,
            "fileUrl": f"uploads/{save_name}",
            "fileType": ext.lstrip(".").upper() if ext else "FILE",
            "fileSize": format_bytes(size),
            "fileSha256": sha256_hex,
            "fileStored": True,
            "preset": False,
            "scoringOrgs": scoring_orgs,
            "roadshowScheduleId": roadshow_schedule_id,
            "roadshowMatchedBy": roadshow_match_method,
            "roadshowMatchedAt": now_iso if roadshow_schedule_id else "",
        }
        store.add_report(report)
        try:
            # 新报告在上传完成时立即进入本地向量索引，后续检索无需再次解析文件。
            _index_knowledge_report(report)
        except Exception:
            # 索引可由运维脚本幂等补建，不能因单个文档抽取异常阻断上传。
            current_app.logger.exception("报告向量化失败，将在下次检索时重试：%s", report_id)
        created.append(report)

    if not created:
        return json_error("没有有效的文件被上传", 400)
    return jsonify({"reports": [public_report(r) for r in created], "count": len(created)})


@app.route("/api/reports/ai-complete", methods=["POST"])
def api_reports_ai_complete():
    """智能补全：接收单个文件，调用 DeepSeek 生成关键词和摘要。
    失败时返回空结果而非错误，不阻断上传流程。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    f = request.files.get("file")
    if not f or not f.filename:
        return json_error("请选择文件", 400)
    ext = os.path.splitext(f.filename)[1].lower()
    if ext not in ALLOWED_UPLOAD_EXTS:
        return json_error("文件格式不支持", 400)

    tmp_fd, tmp_path = tempfile.mkstemp(suffix=ext, prefix="ai_complete_")
    os.close(tmp_fd)
    try:
        f.save(tmp_path)
        external = str(request.form.get("reportType", "internal")).strip() in ("external", "roadshow")
        result = _ai_complete_tags_summary(tmp_path, f.filename, external=external)
        return jsonify(result)
    except Exception as e:
        # 降级：返回空结果，前端提示手动填写
        return jsonify({"tags": [], "summary": "", "author": "", "institution": "", "error": f"智能补全失败：{e}"})
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


@app.route("/api/reports/<rid>/ai-summary", methods=["GET"])
def api_report_ai_summary(rid):
    """读取单篇报告 AI 摘要缓存（每篇幅版本独立缓存，全团队共享）。"""
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("未找到该报告", 404)
    style = str(request.args.get("style") or "standard")
    if style not in AI_SUMMARY_STYLES:
        return json_error("不支持的摘要版本", 400)
    row = _valid_summary_cache(report, style)
    if not row:
        return jsonify({"available": False, "style": style})
    return jsonify({"available": True, **_summary_payload(row, style, row.get("user_id"))})


@app.route("/api/reports/<rid>/ai-summary", methods=["POST"])
def api_report_ai_summary_generate(rid):
    """生成单篇报告 AI 摘要（SSE 流式）。

    前置校验在进流之前完成，失败返回普通 JSON 错误；校验通过后返回
    text/event-stream，事件格式与知识搜索一致：
    data: {"type": "stage"|"delta"|"done"|"error", ...}\n\n
    缓存命中且非 force 时直接回放缓存不调用模型；生成成功后 UPSERT 入库。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("未找到该报告", 404)
    payload = request.get_json(silent=True) or {}
    style = str(payload.get("style") or "standard")
    if style not in AI_SUMMARY_STYLES:
        return json_error("不支持的摘要版本", 400)
    force = bool(payload.get("force"))
    if not _llm_api_key():
        return json_error("AI 摘要尚未配置大模型 API 密钥", 503)
    label = AI_SUMMARY_STYLES[style]["label"]
    cached = _valid_summary_cache(report, style)
    user_id = user["id"]

    def event(data):
        return "data: " + json.dumps(data, ensure_ascii=False) + "\n\n"

    def generate():
        if cached is not None and not force:
            yield event({"type": "stage", "text": "读取已生成的摘要…"})
            yield event({"type": "delta", "text": cached["content"]})
            yield event({"type": "done", **_summary_payload(cached, style, cached.get("user_id"))})
            return
        file_path = report_file_path(report)
        text = _extract_text(file_path, max_chars=AI_SUMMARY_TEXT_CHARS) if file_path else ""
        if not text.strip():
            yield event({"type": "error", "message": "无法从该报告文件提取正文文本，暂不能生成 AI 摘要"})
            return
        yield event({"type": "stage", "text": f"正在生成{label}摘要…"})
        system, prompt = _ai_summary_messages(report, text, style)
        used_provider = []
        content = ""
        try:
            for delta in _stream_llm(prompt, system=system,
                                     max_tokens=AI_SUMMARY_STYLES[style]["max_tokens"],
                                     provider_sink=used_provider):
                content += delta
                yield event({"type": "delta", "text": delta})
            content = content.strip()
            if not content:
                yield event({"type": "error", "message": "模型未返回有效摘要，请稍后重试"})
                return
            model = used_provider[0] if used_provider else ""
            store.save_report_summary(report["id"], style, user_id, content,
                                      model, str(report.get("fileSha256") or ""))
            row = store.get_report_summary(report["id"], style) or {}
            yield event({"type": "done", **_summary_payload(row, style, user_id)})
        except Exception as exc:
            # 不把底层 WinError/代理地址直接暴露给用户，保持提示可读。
            yield event({"type": "error", "message": f"AI 摘要暂时不可用：{exc}"})

    response = Response(stream_with_context(generate()), mimetype="text/event-stream")
    response.headers["Cache-Control"] = "no-cache"
    response.headers["X-Accel-Buffering"] = "no"
    return response


@app.route("/api/reports/<rid>/ask", methods=["POST"])
def api_report_ask(rid):
    """就本文提问：基于该报告已缓存 AI 摘要与全文的单轮问答（SSE 流式）。

    次要功能：不计数、不写库（audit_log 由 after_request 统一记录）。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("未找到该报告", 404)
    question = str((request.get_json(silent=True) or {}).get("question") or "").strip()
    if len(question) < 2:
        return json_error("请输入具体问题", 400)
    if len(question) > 500:
        return json_error("问题请控制在 500 字以内", 400)
    if not _llm_api_key():
        return json_error("问答尚未配置大模型 API 密钥", 503)

    def event(data):
        return "data: " + json.dumps(data, ensure_ascii=False) + "\n\n"

    def generate():
        system, prompt = _report_ask_messages(report, question)
        yield event({"type": "stage", "text": "正在阅读本文并思考…"})
        answer = ""
        try:
            for delta in _stream_llm(prompt, system=system, max_tokens=1500):
                answer += delta
                yield event({"type": "delta", "text": delta})
            yield event({"type": "done", "answer": answer.strip()})
        except Exception as exc:
            yield event({"type": "error", "message": f"问答暂时不可用：{exc}"})

    response = Response(stream_with_context(generate()), mimetype="text/event-stream")
    response.headers["Cache-Control"] = "no-cache"
    response.headers["X-Accel-Buffering"] = "no"
    return response


@app.route("/api/reports/<rid>/file", methods=["GET"])
def api_report_file(rid):
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    rel = report.get("fileUrl", "")
    if not rel:
        return json_error("报告无文件", 404)
    full = report_file_path(report)
    if not full:
        return json_error("文件不存在", 404)
    return send_file(full, as_attachment=True, download_name=report.get("fileName") or os.path.basename(full))


@app.route("/api/reports/<rid>", methods=["DELETE"])
def api_report_delete(rid):
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    # 行政可删任意；普通用户仅删自己上传的
    if user.get("role") == "admin":
        pass
    elif report.get("authorId") == user.get("id"):
        pass
    else:
        return json_error("只能将自己上传的报告移入回收站", 403)
    store.trash_reports([rid], user.get("id"))
    return jsonify({"ok": True, "trashed": True})


@app.route("/api/reports/<rid>", methods=["PUT"])
def api_report_update(rid):
    """报告元信息更新：行政可改任意报告，其他角色仅改本人上传的报告。

    所有可编辑者可改：报告日期、摘要、关键词、内部报告二级分类。
    行政额外可改：报告作者（内部报告按 authorId 重置署名，外部报告改 sourceAuthor）。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    is_admin = user.get("role") == "admin"
    if not is_admin and report.get("authorId") != user.get("id"):
        return json_error("只能修改自己上传的报告", 403)

    data = request.get_json(silent=True) or {}
    fields = {}

    # 所有可编辑者可改的公共字段
    if "reportDate" in data:
        fields["reportDate"] = str(data["reportDate"]).strip()
    if "summary" in data:
        fields["summary"] = str(data["summary"] or "").strip()
    if "tags" in data:
        v = data["tags"]
        if isinstance(v, str):
            v = [t.strip() for t in re.split(r"[,，]", v) if t.strip()][:8]
        elif isinstance(v, list):
            v = [str(t).strip() for t in v if str(t).strip()][:8]
        else:
            v = []
        fields["tags"] = v
    # 研究主题：外部/调研/路演报告可在修改时调整（前端仅对非内部报告发送）
    if "theme" in data:
        theme = str(data["theme"]).strip()
        if theme not in REPORT_THEMES:
            return json_error("研究主题无效", 400)
        fields["theme"] = theme
    # 内部报告二级分类：周报/月报/深度报告/其他报告。月报打分已关闭，
    # 仅改为深度报告时需要打分部门（未配置时默认两部门均打分）。
    if report.get("reportType") == "internal" and "category" in data:
        category = str(data["category"]).strip()
        if category not in ("weekly", "monthly", "deep", "other"):
            return json_error("分类无效", 400)
        fields["category"] = category
        if category == "deep":
            existing = report.get("scoringOrgs")
            if not isinstance(existing, list) or not existing:
                fields["scoringOrgs"] = list(SCORING_ORGS)

    # 仅行政可改报告作者：外部/路演报告改报告作者（sourceAuthor），其余按署名作者处理
    if is_admin:
        if report.get("reportType") in ("external", "roadshow"):
            if "sourceAuthor" in data:
                fields["sourceAuthor"] = str(data["sourceAuthor"] or "").strip()
        else:
            if "authorId" in data:
                author_id = str(data["authorId"]).strip()
                author = store.get_user(author_id)
                if not author:
                    return json_error("请选择有效的报告作者", 400)
                fields["author"] = author["name"]
                fields["authorId"] = author["id"]

    if not fields:
        return json_error("没有需要更新的字段", 400)
    updated = store.update_report(rid, fields)
    return jsonify({"ok": True, "report": public_report(updated)})


@app.route("/api/reports/<rid>/like", methods=["POST"])
def api_report_like(rid):
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    if report.get("reportType", "internal") == "internal":
        return json_error("内部报告不支持点赞", 400)
    return jsonify({"ok": True, **store.toggle_like(rid, user["id"])})


@app.route("/api/reports/<rid>/view", methods=["POST"])
def api_report_view(rid):
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    return jsonify({"ok": True, **store.add_view(rid, user["id"])})


@app.route("/api/reports/<rid>/favorite", methods=["POST"])
def api_report_favorite(rid):
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    if not store.get_report(rid):
        return json_error("报告不存在", 404)
    return jsonify({"ok": True, **store.toggle_favorite(rid, user["id"])})


# --------------------------------------------------------------------------- #
# 路演安排表（调研报告页上方）
# --------------------------------------------------------------------------- #
ROADSHOW_FORMATS = {"online": "线上", "offline": "线下", "hybrid": "线上+线下"}


def _roadshow_week_bounds(anchor_str):
    """返回锚点日期所在周（周一~周日）的起止日期字符串。"""
    try:
        anchor = datetime.strptime(anchor_str, "%Y-%m-%d").date()
    except (TypeError, ValueError):
        anchor = datetime.now(CST).date()
    monday = anchor - timedelta(days=anchor.weekday())
    sunday = monday + timedelta(days=6)
    return monday.strftime("%Y-%m-%d"), sunday.strftime("%Y-%m-%d")


def _validate_roadshow_fields(data):
    """校验路演字段，返回 (event_time, end_time, fmt, institution, room, tencent_id, presenter, topic) 或错误串。"""
    event_time = str(data.get("eventTime", "")).strip()
    try:
        parsed = datetime.strptime(event_time, "%Y-%m-%dT%H:%M")
        event_time = parsed.strftime("%Y-%m-%dT%H:%M")
    except ValueError:
        return "请选择有效的路演时间"
    # 结束时间可选：为空时前端按 1 小时展示；填写时必须晚于开始时间且为同一天
    end_time = str(data.get("endTime", "")).strip()
    if end_time:
        try:
            parsed_end = datetime.strptime(end_time, "%Y-%m-%dT%H:%M")
            end_time = parsed_end.strftime("%Y-%m-%dT%H:%M")
        except ValueError:
            return "请选择有效的结束时间"
        if parsed_end.date() != parsed.date():
            return "结束时间需与路演时间为同一天"
        if parsed_end <= parsed:
            return "结束时间需晚于路演开始时间"
    fmt = str(data.get("format", "")).strip()
    if fmt not in ROADSHOW_FORMATS:
        return "请选择路演形式（线上/线下/线上+线下）"
    institution = str(data.get("institution", "")).strip()[:80]
    organizer = str(data.get("organizer", "")).strip()[:60]
    room = str(data.get("meetingRoom", "")).strip()[:60]
    tencent_id = str(data.get("tencentMeetingId", "")).strip()[:40]
    presenter = str(data.get("presenter", "")).strip()[:60]
    topic = str(data.get("topic", "")).strip()[:200]
    if fmt in ("offline", "hybrid") and not room:
        return "线下/混合路演请填写会议室或地点"
    if fmt in ("online", "hybrid") and not tencent_id:
        return "线上/混合路演请填写腾讯会议号"
    if not presenter:
        return "请填写路演人"
    if not topic:
        return "请填写路演主题"
    return event_time, end_time, fmt, institution, organizer, room, tencent_id, presenter, topic


@app.route("/api/roadshow-schedule", methods=["GET"])
def api_roadshow_schedule_list():
    """按周（周一~周日）返回路演安排；不传 week 参数时返回当前周。

    每条安排附带 reportId/reportIds：该路演已关联的路演报告 id
    （reportId 取最新一篇），供前端区分"已归档/未归档"并提供查看/下载入口。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    week_start, week_end = _roadshow_week_bounds(request.args.get("week", ""))
    items = store.roadshow_items(week_start, week_end)
    report_by_schedule = {}
    for report in sorted(store.reports(), key=lambda r: str(r.get("uploadedAt") or "")):
        schedule_id = str(report.get("roadshowScheduleId") or "")
        if schedule_id:
            report_by_schedule.setdefault(schedule_id, []).append(report.get("id", ""))
    for item in items:
        item["formatLabel"] = ROADSHOW_FORMATS.get(item["format"], item["format"])
        report_ids = report_by_schedule.get(item["id"], [])
        item["reportId"] = report_ids[-1] if report_ids else ""
        item["reportIds"] = report_ids
    return jsonify({"items": items, "weekStart": week_start, "weekEnd": week_end})


ROADSHOW_CAL_START_MIN = 8 * 60   # 周历视图网格 08:00–20:00（与前端一致）
ROADSHOW_CAL_END_MIN = 20 * 60
ROADSHOW_CAL_SLOT_MIN = 30
# 周历视图配色（与前端色块一致）：线上蓝 / 线下绿 / 线上+线下紫
ROADSHOW_CAL_STYLES = {
    "online": {"fill": "EFF6FF", "accent": "2563EB"},
    "offline": {"fill": "ECFDF5", "accent": "059669"},
    "hybrid": {"fill": "F3E8FF", "accent": "7C3AED"},
}


def _roadshow_cal_minutes(value):
    """'2026-08-21T09:30' → 570；解析失败返回 None。"""
    text = str(value or "")
    try:
        return int(text[11:13]) * 60 + int(text[14:16])
    except (ValueError, IndexError):
        return None


def _append_roadshow_calendar_sheet(wb, items, week_start):
    """在导出工作簿追加"周历视图" sheet：布局仿前端路演安排界面——
    周一~周五 × 8:00–20:00 半小时格，重叠路演并排分栏（不重叠的占整行宽），
    色块按形式配色并带左侧色条，单元格合并后可直接编辑文字。"""
    ws = wb.create_sheet("周历视图")
    thin = Side(style="thin", color="D1D5DB")
    grid_border = Border(left=thin, right=thin, top=thin, bottom=thin)
    slot_count = (ROADSHOW_CAL_END_MIN - ROADSHOW_CAL_START_MIN) // ROADSHOW_CAL_SLOT_MIN
    header_row, first_slot_row = 1, 2

    start_date = datetime.strptime(week_start, "%Y-%m-%d").date()
    weekdays = "一二三四五六日"
    day_items, weekend_items = [[] for _ in range(5)], []
    for item in items:
        try:
            day = datetime.strptime(str(item.get("event_time", ""))[:10], "%Y-%m-%d").date()
        except ValueError:
            continue
        offset = (day - start_date).days
        if 0 <= offset < 5:
            day_items[offset].append(item)
        elif 0 <= offset < 7:
            weekend_items.append(item)

    # 每天先算重叠分栏（与前端聚簇逻辑一致）：返回 [{item, track, tracks}] 及该天栏数
    def day_blocks(day_list):
        blocks = []
        for item in day_list:
            start = _roadshow_cal_minutes(item.get("event_time"))
            if start is None:
                continue
            end = _roadshow_cal_minutes(item.get("end_time"))
            end = end if end and end > start else start + 60  # 未填结束时间默认 1 小时
            blocks.append({"item": item,
                           "start": max(start, ROADSHOW_CAL_START_MIN),
                           "end": min(end, ROADSHOW_CAL_END_MIN)})
        blocks.sort(key=lambda b: (b["start"], -b["end"]))
        clusters, tracks = [], 1
        for block in blocks:
            if clusters and block["start"] < clusters[-1]["end"]:
                clusters[-1]["blocks"].append(block)
                clusters[-1]["end"] = max(clusters[-1]["end"], block["end"])
            else:
                clusters.append({"blocks": [block], "end": block["end"]})
        for cluster in clusters:
            track_ends = []
            for block in cluster["blocks"]:
                idx = next((i for i, end in enumerate(track_ends) if end <= block["start"]), None)
                if idx is None:
                    track_ends.append(block["end"])
                    idx = len(track_ends) - 1
                else:
                    track_ends[idx] = block["end"]
                block["track"] = idx
            for block in cluster["blocks"]:
                block["tracks"] = len(track_ends)
            tracks = max(tracks, len(track_ends))
        return blocks, tracks

    layout = [day_blocks(day_list) for day_list in day_items]

    # 表头与时间轴
    header_font = Font(bold=True, size=11)
    center = Alignment(horizontal="center", vertical="center", wrap_text=True)
    ws.cell(row=header_row, column=1, value="时间").font = header_font
    ws.cell(row=header_row, column=1).alignment = center
    ws.cell(row=header_row, column=1).border = grid_border
    for slot in range(slot_count):
        row = first_slot_row + slot
        minutes = ROADSHOW_CAL_START_MIN + slot * ROADSHOW_CAL_SLOT_MIN
        label = ws.cell(row=row, column=1, value=f"{minutes // 60:02d}:{minutes % 60:02d}"
                        if minutes % 60 == 0 else None)
        label.border = grid_border
        label.font = Font(size=9, color="6B7280")
        label.alignment = Alignment(horizontal="right", vertical="top")
        ws.row_dimensions[row].height = 20
    ws.column_dimensions["A"].width = 7
    ws.row_dimensions[header_row].height = 32

    # 周一~周五列：按各天栏数切分，铺网格底纹后放置路演色块
    left_align = Alignment(horizontal="left", vertical="top", wrap_text=True)
    col = 2
    for offset, (blocks, tracks) in enumerate(layout):
        day = start_date + timedelta(days=offset)
        for t in range(tracks):
            ws.column_dimensions[get_column_letter(col + t)].width = 26 if tracks == 1 else 13
        header = ws.cell(row=header_row, column=col,
                         value=f"周{weekdays[offset]}\n{day.month}.{day.day}")
        header.font = header_font
        header.alignment = center
        for t in range(tracks):
            ws.cell(row=header_row, column=col + t).border = grid_border
        if tracks > 1:
            ws.merge_cells(start_row=header_row, start_column=col,
                           end_row=header_row, end_column=col + tracks - 1)
        for row in range(first_slot_row, first_slot_row + slot_count):
            for t in range(tracks):
                ws.cell(row=row, column=col + t).border = grid_border
        for block in blocks:
            item = block["item"]
            # 横向：按所在簇的栏数占当天宽度的等分（栏数少于当天总栏数时合并多列占满）
            c0 = col + block["track"] * tracks // block["tracks"]
            c1 = col + (block["track"] + 1) * tracks // block["tracks"] - 1
            # 纵向：跨半小时行合并
            r0 = first_slot_row + (block["start"] - ROADSHOW_CAL_START_MIN) // ROADSHOW_CAL_SLOT_MIN
            r1 = first_slot_row + (block["end"] - ROADSHOW_CAL_START_MIN) // ROADSHOW_CAL_SLOT_MIN - 1
            style = ROADSHOW_CAL_STYLES.get(item.get("format"), {"fill": "F3F4F6", "accent": "6B7280"})
            meta = item.get("format") in ("online", "hybrid") and item.get("tencent_meeting_id") \
                or item.get("meeting_room") or ""
            institution = item.get("institution", "")
            organizer_suffix = f"（主约:{item['organizer']}）" if item.get("organizer") else ""
            lines = [item.get("topic") or "未填主题",
                     f"{block['start'] // 60:02d}:{block['start'] % 60:02d}"
                     f"-{block['end'] // 60:02d}:{block['end'] % 60:02d} · {item.get('presenter', '')}{organizer_suffix}",
                     " · ".join(part for part in (institution, ROADSHOW_FORMATS.get(item.get("format"), ""), meta) if part)]
            accent = Side(style="thick", color=style["accent"])
            for row in range(r0, r1 + 1):
                for column in range(c0, c1 + 1):
                    cell = ws.cell(row=row, column=column)
                    cell.fill = PatternFill("solid", fgColor=style["fill"])
                    cell.border = Border(left=accent if column == c0 else thin,
                                         right=thin, top=thin, bottom=thin)
            cell = ws.cell(row=r0, column=c0, value="\n".join(lines))
            cell.font = Font(size=9)
            cell.alignment = left_align
            ws.merge_cells(start_row=r0, start_column=c0, end_row=r1, end_column=c1)
        col += tracks

    ws.freeze_panes = "B2"

    # 周末路演：仿前端，在网格下方列出
    if weekend_items:
        note_row = first_slot_row + slot_count + 1
        parts = [f"{str(it.get('event_time', ''))[5:16].replace('T', ' ')} {it.get('presenter', '')}《{it.get('topic', '')}》"
                 for it in weekend_items]
        ws.merge_cells(start_row=note_row, start_column=1,
                       end_row=note_row, end_column=max(col - 1, 6))
        note = ws.cell(row=note_row, column=1,
                       value=f"周末另有 {len(weekend_items)} 场路演：" + "；".join(parts))
        note.alignment = Alignment(horizontal="left", vertical="center", wrap_text=True)
        note.font = Font(size=10, color="92400E")
        note.fill = PatternFill("solid", fgColor="FFFBEB")


@app.route("/api/roadshow-schedule/export", methods=["GET"])
def api_roadshow_schedule_export():
    """导出一周路演安排 Excel（仅行政角色），week 参数与列表接口一致。"""
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    if user.get("role") != "admin":
        return json_error("仅行政角色可导出", 403)
    week_start, week_end = _roadshow_week_bounds(request.args.get("week", ""))
    items = store.roadshow_items(week_start, week_end)
    weekdays = "一二三四五六日"

    def hm(value):
        return str(value or "")[11:16]

    wb = Workbook()
    ws = wb.active
    ws.title = "路演安排"
    ws.append(["日期", "星期", "开始时间", "结束时间", "形式", "路演机构", "主约人",
               "路演人", "路演主题", "会议室/地点", "腾讯会议号", "登记人"])
    for item in items:
        event_time = str(item.get("event_time", ""))
        try:
            day = datetime.strptime(event_time[:10], "%Y-%m-%d").date()
            weekday = f"周{weekdays[day.weekday()]}"
        except ValueError:
            weekday = ""
        ws.append([event_time[:10], weekday, hm(event_time), hm(item.get("end_time")),
                   ROADSHOW_FORMATS.get(item.get("format"), str(item.get("format", ""))),
                   item.get("institution", ""), item.get("organizer", ""), item.get("presenter", ""),
                   item.get("topic", ""), item.get("meeting_room", ""), item.get("tencent_meeting_id", ""),
                   item.get("created_by_name", "")])
    for cell in ws[1]:
        cell.font = Font(bold=True)
        cell.fill = PatternFill("solid", fgColor="EEF2FF")
        cell.alignment = Alignment(horizontal="center", vertical="center")
    for idx, width in enumerate([12, 7, 10, 10, 11, 16, 10, 10, 50, 14, 16, 10], start=1):
        ws.column_dimensions[get_column_letter(idx)].width = width
    ws.freeze_panes = "A2"
    _append_roadshow_calendar_sheet(wb, items, week_start)
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return send_file(buf, as_attachment=True,
                     download_name=f"路演安排_{week_start}_{week_end}.xlsx",
                     mimetype="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet")


@app.route("/api/roadshow-schedule", methods=["POST"])
def api_roadshow_schedule_add():
    """新增路演安排：所有登录用户可新增，创建人=本人。"""
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    data = request.get_json(silent=True) or {}
    validated = _validate_roadshow_fields(data)
    if isinstance(validated, str):
        return json_error(validated, 400)
    event_time, end_time, fmt, institution, organizer, room, tencent_id, presenter, topic = validated
    # 主约人默认为账户本人；行政可在前端选择其他成员代为登记
    organizer = organizer or user["name"]
    item = {
        "id": f"roadshow-{int(time.time() * 1000)}-{uuid.uuid4().hex[:8]}",
        "eventTime": event_time,
        "endTime": end_time,
        "format": fmt,
        "institution": institution,
        "organizer": organizer,
        "meetingRoom": room,
        "tencentMeetingId": tencent_id,
        "presenter": presenter,
        "topic": topic,
        "createdById": user["id"],
        "createdByName": user["name"],
    }
    store.add_roadshow_item(item)
    return jsonify({"ok": True, "item": {**item, "formatLabel": ROADSHOW_FORMATS[fmt]}})


@app.route("/api/roadshow-schedule/<item_id>", methods=["PUT"])
def api_roadshow_schedule_update(item_id):
    """修改路演安排：创建人可改自己的记录，行政可改任何人的记录。"""
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    item = store.get_roadshow_item(item_id)
    if not item:
        return json_error("路演安排不存在", 404)
    if item.get("created_by") != user["id"] and user.get("role") != "admin":
        return json_error("只能修改自己创建的路演安排", 403)
    data = request.get_json(silent=True) or {}
    validated = _validate_roadshow_fields(data)
    if isinstance(validated, str):
        return json_error(validated, 400)
    event_time, end_time, fmt, institution, organizer, room, tencent_id, presenter, topic = validated
    # 主约人为空时保留原值（行政代登记后由他人补充修改的场景）
    organizer = organizer or item.get("organizer") or user["name"]
    updated = store.update_roadshow_item(item_id, {
        "eventTime": event_time, "endTime": end_time, "format": fmt,
        "institution": institution, "organizer": organizer,
        "meetingRoom": room, "tencentMeetingId": tencent_id,
        "presenter": presenter, "topic": topic,
    })
    updated["formatLabel"] = ROADSHOW_FORMATS.get(updated["format"], updated["format"])
    return jsonify({"ok": True, "item": updated})


@app.route("/api/roadshow-schedule/<item_id>", methods=["DELETE"])
def api_roadshow_schedule_delete(item_id):
    """删除路演安排：创建人可删自己的记录，行政可删任何人的记录。"""
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    item = store.get_roadshow_item(item_id)
    if not item:
        return json_error("路演安排不存在", 404)
    if item.get("created_by") != user["id"] and user.get("role") != "admin":
        return json_error("只能删除自己创建的路演安排", 403)
    store.delete_roadshow_item(item_id)
    return jsonify({"ok": True})


# 路演机构常见简称 → 全称映射（识别提示词用，命中即返回全称）
ROADSHOW_INSTITUTION_ALIASES = {
    "兴证": "兴业证券", "中金": "中金公司", "中信": "中信证券", "广发": "广发证券",
    "国君": "国泰君安证券", "国泰君安": "国泰君安证券", "海通": "海通证券",
    "华泰": "华泰证券", "招商证券": "招商证券", "招商固收": "招商证券", "招商": "招商证券",
    "申万": "申万宏源证券", "申万宏源": "申万宏源证券",
    "中信建投": "中信建投证券", "建投": "中信建投证券", "国信": "国信证券",
    "东方": "东方证券", "光大": "光大证券", "浙商": "浙商证券", "民生": "民生证券",
    "安信": "国投证券", "天风": "天风证券", "兴业研究": "兴业研究",
    "国盛": "国盛证券", "华西": "华西证券", "东吴": "东吴证券", "国金": "国金证券",
    "长江": "长江证券", "开源": "开源证券", "信达": "信达证券", "华福": "华福证券",
    "易方达": "易方达基金", "华夏": "华夏基金", "嘉实": "嘉实基金", "南方基金": "南方基金",
    "博时": "博时基金", "富国": "富国基金", "中欧": "中欧基金", "汇添富": "汇添富基金",
    "鹏华": "鹏华基金", "兴全": "兴证全球基金",
}


def _roadshow_ai_parse_prompt(text):
    aliases = "、".join(f"{short}→{full}" for short, full in ROADSHOW_INSTITUTION_ALIASES.items())
    return (
        "你是路演信息抽取助手。请从下面的路演通知文本中提取字段，只输出一个JSON对象：\n"
        '{"date": "…", "weekday": "…", "time": "…", "format": "…", "institution": "…", '
        '"tencentMeetingId": "…", "meetingRoom": "…", "presenter": "…", "topic": "…"}\n\n'
        "规则：\n"
        "1. date：路演日期。文本含年份时输出 YYYY-MM-DD；没有年份（如“9月25日”）时只输出 MM-DD；"
        "只有星期（如“周四”）没有具体日期时 date 输出空字符串，并把星期写到 weekday（1=周一 … 5=周五，无则空字符串）\n"
        "2. time：开始时间 HH:MM（24小时制），没有则输出空字符串\n"
        "3. format：线上路演输出 online，线下路演输出 offline，同时包含线上和线下输出 hybrid；"
        "无法判断时，含腾讯会议号输出 online，含会议室/地点输出 offline\n"
        "4. institution：路演机构。文本中的机构简称要换成全称，常见对照：" + aliases + "；"
        "不在对照中的简称按常识补全为全称；无法确定输出空字符串\n"
        "5. tencentMeetingId：腾讯会议号，如 659-689-968，保留原样；没有输出空字符串\n"
        "6. meetingRoom：会议室或地点，如 9层一会；没有输出空字符串\n"
        "7. presenter：路演人姓名；没有输出空字符串\n"
        "8. topic：路演主题；没有输出空字符串\n"
        "9. 严格基于文本内容，不要编造；没有的字段一律输出空字符串\n\n"
        f"路演通知文本：\n{text}"
    )


def _nearest_future_date(month, day):
    """无年份的月日按最近的未来日期补全年份（如 8 月看到“9月25日”→ 今年 9 月 25 日）。"""
    now = datetime.now(CST)
    for year in (now.year, now.year + 1):
        try:
            candidate = datetime(year, month, day)
        except ValueError:
            return None
        if candidate.date() >= now.date():
            return candidate.strftime("%Y-%m-%d")
    return None


@app.route("/api/roadshow-schedule/ai-parse", methods=["POST"])
def api_roadshow_schedule_ai_parse():
    """粘贴路演通知文本，调用大模型（MiMo 默认，DeepSeek 兜底）识别路演要素。

    日期补全：完整日期直接用；只有月日按最近的未来日期补全年份；只有星期几时
    按 weekStart（当前显示周的周一）推算，默认落在正在查看的这一周。
    无结束时间由前端默认按 1 小时处理。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    data = request.get_json(silent=True) or {}
    text = str(data.get("text", "")).strip()
    week_start = str(data.get("weekStart", "")).strip()
    if not text:
        return json_error("请粘贴路演通知文本", 400)
    if len(text) > 2000:
        return json_error("文本请控制在 2000 字以内", 400)
    try:
        raw = _call_llm(_roadshow_ai_parse_prompt(text))
        parsed = json.loads(raw)
    except Exception as exc:
        return json_error(f"自动识别失败：{exc}", 503)

    # 日期补全：完整日期 > 月日（最近未来）> 星期几（按显示周推算）
    date_str = str(parsed.get("date", "")).strip()
    if re.fullmatch(r"\d{1,2}-\d{1,2}", date_str):
        month, day = (int(part) for part in date_str.split("-"))
        date_str = _nearest_future_date(month, day) or ""
    if not re.fullmatch(r"\d{4}-\d{2}-\d{2}", date_str):
        date_str = ""
    if not date_str and week_start:
        try:
            weekday = int(str(parsed.get("weekday", "")).strip())
            monday = datetime.strptime(week_start, "%Y-%m-%d")
            if 1 <= weekday <= 5:
                date_str = (monday + timedelta(days=weekday - 1)).strftime("%Y-%m-%d")
        except (TypeError, ValueError):
            pass
    time_str = str(parsed.get("time", "")).strip()
    if not re.fullmatch(r"\d{1,2}:\d{2}", time_str):
        time_str = ""
    fmt = str(parsed.get("format", "")).strip()
    if fmt not in ROADSHOW_FORMATS:
        fmt = ""
    result = {
        "eventTime": f"{date_str}T{time_str}" if date_str and time_str else (date_str or ""),
        "format": fmt,
        "institution": str(parsed.get("institution", "")).strip()[:80],
        "tencentMeetingId": str(parsed.get("tencentMeetingId", "")).strip()[:40],
        "meetingRoom": str(parsed.get("meetingRoom", "")).strip()[:60],
        "presenter": str(parsed.get("presenter", "")).strip()[:60],
        "topic": str(parsed.get("topic", "")).strip()[:200],
    }
    return jsonify(result)


# --------------------------------------------------------------------------- #
# 路演报告 ↔ 路演安排匹配：规则 + 大模型自动匹配，手工匹配兜底
# --------------------------------------------------------------------------- #
ROADSHOW_MATCH_WINDOW_DAYS = 3   # 自动匹配的日期窗口（报告日期前后各 N 天）
ROADSHOW_MATCH_RULE_SCORE = 5    # 规则匹配置信分阈值（日期3+路演人3+机构2+主题2）
ROADSHOW_MATCH_RULE_MARGIN = 2   # 规则匹配需领先第二名的分差，避免多场近似时误配


def _normalize_roadshow_institution(name):
    """机构简称归一化：命中常见简称映射时换全称，方便与报告机构比对。"""
    text = str(name or "").strip()
    return ROADSHOW_INSTITUTION_ALIASES.get(text, text)


ROADSHOW_MATCH_TOKEN_SPLIT = re.compile(r"[\s,，。；;、：:（）()《》\"'?!？\-—/\\]+")


def _roadshow_match_score(event_time, presenter, institution, topic,
                          report_date, source_author, source_institution, title):
    """报告与路演安排的匹配打分，两个方向的自动匹配共用。

    日期同日 +3（差一天 +1）；路演人与报告作者互含 +3；机构归一化后互含 +2；
    主题与标题 token 重合最多 +2。满分 10。
    """
    value = 0
    try:
        day = datetime.strptime(str(event_time or "")[:10], "%Y-%m-%d").date()
    except ValueError:
        day = None
    try:
        anchor = datetime.strptime(str(report_date or ""), "%Y-%m-%d").date()
    except ValueError:
        anchor = None
    if day is not None and anchor is not None:
        if day == anchor:
            value += 3
        elif abs((day - anchor).days) <= 1:
            value += 1
    presenter = str(presenter or "")
    author = str(source_author or "")
    if presenter and author and (presenter in author or author in presenter):
        value += 3
    schedule_inst = _normalize_roadshow_institution(institution)
    report_inst = _normalize_roadshow_institution(source_institution)
    if schedule_inst and report_inst and (schedule_inst in report_inst or report_inst in schedule_inst):
        value += 2
    title_tokens = {t for t in ROADSHOW_MATCH_TOKEN_SPLIT.split(str(title or "")) if len(t) >= 2}
    topic_tokens = {t for t in ROADSHOW_MATCH_TOKEN_SPLIT.split(str(topic or "")) if len(t) >= 2}
    value += min(2, len(title_tokens & topic_tokens))
    return value


def _roadshow_auto_match(report_date, source_author, source_institution, title):
    """上传路演报告时自动匹配路演安排。

    规则优先（日期/路演人/机构/主题重合打分），规则无把握时把候选交给
    大模型挑选。返回 (schedule_id, method)，method 为 rule/llm；
    匹配不上返回 ("", "")。任何异常都不阻断上传。
    """
    try:
        anchor = datetime.strptime(str(report_date or ""), "%Y-%m-%d").date()
    except ValueError:
        return "", ""
    date_from = (anchor - timedelta(days=ROADSHOW_MATCH_WINDOW_DAYS)).strftime("%Y-%m-%d")
    date_to = (anchor + timedelta(days=ROADSHOW_MATCH_WINDOW_DAYS)).strftime("%Y-%m-%d")
    items = store.roadshow_items(date_from, date_to)
    if not items:
        return "", ""

    def score(item):
        return _roadshow_match_score(item.get("event_time"), item.get("presenter"),
                                     item.get("institution"), item.get("topic"),
                                     report_date, source_author, source_institution, title)

    ranked = sorted(items, key=score, reverse=True)
    best_score = score(ranked[0])
    runner_up = score(ranked[1]) if len(ranked) > 1 else -1
    if best_score >= ROADSHOW_MATCH_RULE_SCORE and best_score - runner_up >= ROADSHOW_MATCH_RULE_MARGIN:
        return ranked[0]["id"], "rule"

    # 规则无把握：把候选交给大模型挑选（失败/超时则留空，用户可手工匹配）
    try:
        lines = []
        for idx, item in enumerate(ranked, 1):
            when = str(item.get("event_time", ""))[:16].replace("T", " ")
            lines.append(f"{idx}. id={item['id']} | {when} | 路演人:{item.get('presenter', '')} | "
                         f"机构:{item.get('institution', '')} | 主题:{item.get('topic', '')}")
        prompt = (
            "你是路演报告归档助手。一篇路演报告的信息如下：\n"
            f"报告日期：{report_date}\n报告标题：{title}\n"
            f"报告作者（路演人）：{source_author or '未知'}\n报告机构：{source_institution or '未知'}\n\n"
            "候选路演安排：\n" + "\n".join(lines) + "\n\n"
            "请判断哪一场路演最可能是这篇报告对应的路演。只输出一个JSON对象："
            '{"scheduleId": "候选id或空字符串"}\n'
            "规则：选信息最吻合的一场（日期、路演人、机构、主题）；都不吻合时 scheduleId 输出空字符串；"
            "严格基于给定信息判断，不要编造。"
        )
        raw = _call_llm(prompt)
        chosen = str(json.loads(raw).get("scheduleId", "")).strip()
        if chosen and chosen in {item["id"] for item in items}:
            return chosen, "llm"
    except Exception as exc:
        current_app.logger.warning("路演报告大模型匹配失败：%s", exc)
    return "", ""


@app.route("/api/roadshow-schedule/options", methods=["GET"])
def api_roadshow_schedule_options():
    """手工匹配用：返回锚点日期前后各 N 天（默认 10 天）的路演安排候选项。"""
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    try:
        anchor = datetime.strptime(str(request.args.get("date", "")), "%Y-%m-%d").date()
    except ValueError:
        anchor = datetime.now(CST).date()
    try:
        days = min(max(int(request.args.get("days", 10)), 1), 60)
    except (TypeError, ValueError):
        days = 10
    date_from = (anchor - timedelta(days=days)).strftime("%Y-%m-%d")
    date_to = (anchor + timedelta(days=days)).strftime("%Y-%m-%d")
    options = []
    for item in store.roadshow_items(date_from, date_to):
        when = str(item.get("event_time", ""))[:16].replace("T", " ")
        end = str(item.get("end_time") or "")[11:16]
        label = f"{when}{f'-{end}' if end else ''} · {item.get('presenter', '')}《{item.get('topic', '')}》"
        if item.get("institution"):
            label += f"（{item['institution']}）"
        options.append({"id": item["id"], "label": label})
    return jsonify({"options": options})


@app.route("/api/roadshow-schedule/match", methods=["POST"])
def api_roadshow_schedule_match():
    """手工建立/解除路演报告与路演安排的关联。

    body: {reportId, scheduleId}；scheduleId 传空字符串表示取消关联。
    权限：行政，或相关本人（报告署名作者/实际上传人/目标或现有路演安排的创建人）。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    data = request.get_json(silent=True) or {}
    report_id = str(data.get("reportId", "")).strip()
    schedule_id = str(data.get("scheduleId", "")).strip()
    report = store.get_report(report_id)
    if not report:
        return json_error("报告不存在", 404)
    if report.get("reportType") != "roadshow":
        return json_error("仅路演报告支持关联路演安排", 400)
    schedule = None
    if schedule_id:
        schedule = store.get_roadshow_item(schedule_id)
        if not schedule:
            return json_error("路演安排不存在，请刷新后重试", 404)
    current_id = str(report.get("roadshowScheduleId") or "")
    current_schedule = store.get_roadshow_item(current_id) if current_id else None
    involved = (report.get("authorId") == user.get("id")
                or report.get("uploadedById") == user.get("id")
                or (schedule is not None and schedule.get("created_by") == user.get("id"))
                or (current_schedule is not None and current_schedule.get("created_by") == user.get("id")))
    if user.get("role") != "admin" and not involved:
        return json_error("仅行政或相关本人可匹配路演报告", 403)
    fields = {"roadshowScheduleId": schedule_id,
              "roadshowMatchedBy": "manual" if schedule_id else "",
              "roadshowMatchedAt": _now_iso() if schedule_id else ""}
    updated = store.update_report(report_id, fields)
    return jsonify({"ok": True, "report": public_report(updated)})


@app.route("/api/roadshow-schedule/<item_id>/auto-match-report", methods=["POST"])
def api_roadshow_auto_match_report(item_id):
    """为路演安排推荐一份最吻合的未关联路演报告（规则 + 大模型，报告侧的反向入口）。

    只推荐不写库：返回 {recommended, method, report, message}，前端展示推荐结果，
    由人工确认后再调用 /api/roadshow-schedule/match 保存关联。
    权限：行政或该路演安排创建人。
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    item = store.get_roadshow_item(item_id)
    if not item:
        return json_error("路演安排不存在", 404)
    if user.get("role") != "admin" and item.get("created_by") != user.get("id"):
        return json_error("仅行政或路演创建人可自动匹配报告", 403)
    try:
        day = datetime.strptime(str(item.get("event_time", ""))[:10], "%Y-%m-%d").date()
    except ValueError:
        return json_error("该路演安排缺少有效日期，无法自动匹配", 400)
    date_from = (day - timedelta(days=ROADSHOW_MATCH_WINDOW_DAYS)).strftime("%Y-%m-%d")
    date_to = (day + timedelta(days=ROADSHOW_MATCH_WINDOW_DAYS)).strftime("%Y-%m-%d")
    candidates = [
        report for report in store.reports()
        if report.get("reportType") == "roadshow"
        and not str(report.get("roadshowScheduleId") or "")
        and date_from <= str(report.get("reportDate") or "") <= date_to
    ]
    if not candidates:
        return jsonify({"recommended": False, "method": "", "report": None,
                        "message": "该时段前后没有未关联的路演报告，可先上传或展开选择"})

    def score(report):
        return _roadshow_match_score(item.get("event_time"), item.get("presenter"),
                                     item.get("institution"), item.get("topic"),
                                     report.get("reportDate"), report.get("sourceAuthor"),
                                     report.get("sourceInstitution"), report.get("title"))

    ranked = sorted(candidates, key=score, reverse=True)
    best_score = score(ranked[0])
    runner_up = score(ranked[1]) if len(ranked) > 1 else -1
    chosen, method = "", ""
    if best_score >= ROADSHOW_MATCH_RULE_SCORE and best_score - runner_up >= ROADSHOW_MATCH_RULE_MARGIN:
        chosen, method = ranked[0]["id"], "rule"
    else:
        # 规则无把握：把候选交给大模型挑选（失败/超时则留空，用户可手工选择）
        try:
            lines = []
            for idx, report in enumerate(ranked, 1):
                lines.append(f"{idx}. reportId={report['id']} | 日期:{report.get('reportDate', '')} | "
                             f"标题:{report.get('title', '')} | 作者:{report.get('sourceAuthor', '')} | "
                             f"机构:{report.get('sourceInstitution', '')}")
            prompt = (
                "你是路演报告归档助手。一场路演安排的信息如下：\n"
                f"路演日期：{item.get('event_time', '')[:10]}\n主题：{item.get('topic', '')}\n"
                f"路演人：{item.get('presenter', '') or '未知'}\n机构：{item.get('institution', '') or '未知'}\n\n"
                "候选未关联的路演报告：\n" + "\n".join(lines) + "\n\n"
                "请判断哪一篇报告最可能是这场路演对应的报告。只输出一个JSON对象："
                '{"reportId": "候选reportId或空字符串"}\n'
                "规则：选信息最吻合的一篇（日期、路演人、机构、主题）；都不吻合时 reportId 输出空字符串；"
                "严格基于给定信息判断，不要编造。"
            )
            raw = _call_llm(prompt)
            picked = str(json.loads(raw).get("reportId", "")).strip()
            if picked and picked in {report["id"] for report in candidates}:
                chosen, method = picked, "llm"
        except Exception as exc:
            current_app.logger.warning("路演安排大模型匹配报告失败：%s", exc)

    if not chosen:
        return jsonify({"recommended": False, "method": "", "report": None,
                        "message": "未找到足够吻合的路演报告，请手工选择或展开已关联报告"})
    report = next(report for report in candidates if report["id"] == chosen)
    return jsonify({"recommended": True, "method": method, "report": public_report(report),
                    "message": ""})


@app.route("/api/work-reminders", methods=["GET"])
def api_work_reminders():
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    config = store.reminder_config()
    period = config.get("period") or datetime.now(CST).strftime("%Y")
    report_category = config.get("reportCategory", "deep")
    year_reports = [
        report for report in store.reports()
        if str(report.get("reportDate") or report.get("uploadedAt") or "")[:4] == period
    ]
    topic_reports = [report for report in year_reports
                     if report.get("reportType") == "internal" and report.get("category") == report_category]
    users = store.users()
    privileged = user.get("role") in ("leader", "admin")
    configured_rules = config.get("rules", [])
    if privileged:
        visible_rules = configured_rules
        visible_user_ids = {item["id"] for item in users}
    else:
        visible_rules = [rule for rule in configured_rules if user["id"] in (rule.get("userIds") or [])]
        visible_user_ids = {user["id"]}
        for rule in visible_rules:
            visible_user_ids.update(rule.get("userIds") or [])
    people = []
    for item in users:
        if item["id"] not in visible_user_ids:
            continue
        uploaded = [report for report in year_reports if report.get("authorId") == item["id"]]
        people.append({
            "id": item["id"], "name": item["name"], "org": item.get("org", ""),
            "count": len(uploaded),
            "reports": [{"id": report["id"], "title": report.get("title", ""),
                         "reportType": report.get("reportType", "internal"),
                         "category": report.get("category", "other")} for report in uploaded],
        })
    rules = []
    users_by_id = {item["id"]: item for item in users}
    for rule in visible_rules:
        member_ids = [uid for uid in rule.get("userIds", []) if uid in users_by_id]
        matched = [report for report in topic_reports if report.get("authorId") in member_ids]
        target = max(int(rule.get("target", 0) or 0), 0)
        rules.append({
            "id": rule.get("id", ""), "label": rule.get("label", "未命名要求"),
            "mode": rule.get("mode", "group"), "target": target,
            "completed": len(matched), "remaining": max(target - len(matched), 0),
            "members": [users_by_id[uid]["name"] for uid in member_ids],
            "reports": [{"id": report["id"], "title": report.get("title", "")} for report in matched],
        })
    return jsonify({
        "period": period, "reportCategory": report_category, "people": people, "rules": rules,
        "summary": {"target": sum(row["target"] for row in rules),
                    "completed": sum(min(row["completed"], row["target"]) for row in rules),
                    "remaining": sum(row["remaining"] for row in rules)},
    })


@app.route("/api/knowledge-search", methods=["GET", "POST", "DELETE"])
def api_knowledge_search():
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    if request.method == "DELETE":
        # 用户清空自己的知识搜索历史，不影响当日已用额度。
        store.clear_qa_history(user["id"])
        return jsonify({"ok": True})
    day = datetime.now(CST).strftime("%Y-%m-%d")
    used = store.qa_usage_today(user["id"], day)
    knowledge_config = store.knowledge_config()
    limit = knowledge_config["leaderLimit"] if user.get("role") == "leader" else knowledge_config["memberLimit"]
    if request.method == "GET":
        history = store.qa_history_for_user(user["id"])
        return jsonify({"limit": limit, "used": used, "remaining": max(limit - used, 0),
                        "available": bool(_llm_api_key()), "history": history})
    if used >= limit:
        return json_error(f"今日 {limit} 次知识搜索额度已用完，请明天再试", 429)
    payload = request.get_json(silent=True) or {}
    question = str(payload.get("question", "")).strip()
    if len(question) < 2:
        return json_error("请输入具体问题", 400)
    if len(question) > 300:
        return json_error("问题请控制在 300 字以内", 400)
    if not _llm_api_key():
        return json_error("知识搜索尚未配置大模型 API 密钥", 503)
    filters = _parse_knowledge_filters(payload)
    try:
        result = _answer_knowledge_question(question, filters=filters)
    except Exception as exc:
        # 不把底层 WinError/代理地址直接暴露给用户，便于定位并保持提示可读。
        return json_error(f"知识搜索暂时不可用：{exc}", 503)
    store.add_qa_usage(user["id"], day, question)
    store.add_qa_history(user["id"], question, result.get("answer", ""), result.get("sources", []))
    return jsonify({**result, "limit": limit, "used": used + 1, "remaining": max(limit - used - 1, 0)})


@app.route("/api/knowledge-search/stream", methods=["POST"])
def api_knowledge_search_stream():
    """知识搜索流式端点：以 SSE 推送检索进度与回答增量，缓解长等待焦虑。

    前置校验（登录/额度/问题长度/密钥）在进入流之前完成，失败照常返回
    JSON 错误；校验通过后返回 text/event-stream，事件格式为
    data: {"type": "stage"|"delta"|"done"|"error", ...}\n\n
    """
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    day = datetime.now(CST).strftime("%Y-%m-%d")
    used = store.qa_usage_today(user["id"], day)
    knowledge_config = store.knowledge_config()
    limit = knowledge_config["leaderLimit"] if user.get("role") == "leader" else knowledge_config["memberLimit"]
    if used >= limit:
        return json_error(f"今日 {limit} 次知识搜索额度已用完，请明天再试", 429)
    payload = request.get_json(silent=True) or {}
    question = str(payload.get("question", "")).strip()
    if len(question) < 2:
        return json_error("请输入具体问题", 400)
    if len(question) > 300:
        return json_error("问题请控制在 300 字以内", 400)
    if not _llm_api_key():
        return json_error("知识搜索尚未配置大模型 API 密钥", 503)
    filters = _parse_knowledge_filters(payload)

    user_id = user["id"]
    remaining = max(limit - used - 1, 0)

    def event(payload):
        return "data: " + json.dumps(payload, ensure_ascii=False) + "\n\n"

    def generate():
        # 阶段一：本地检索。全文抽取是本地主要耗时，逐份报告推送解析进度。
        intent = _knowledge_intent(question)
        candidate_limit = 12 if intent == KNOWLEDGE_INTENT_GENERAL else 6
        candidates = []
        retrieval = _iter_knowledge_candidates(question, limit=candidate_limit, filters=filters)
        while True:
            try:
                done_count, total = next(retrieval)
            except StopIteration as stop:
                candidates = stop.value or []
                break
            yield event({"type": "stage", "text": f"正在检索向量索引（{done_count}/{total}）…"})
        if not candidates:
            answer = _knowledge_no_match_answer(filters)
            store.add_qa_usage(user_id, day, question)
            store.add_qa_history(user_id, question, answer, [])
            yield event({"type": "done", "answer": answer, "sources": [],
                         "limit": limit, "used": used + 1, "remaining": remaining})
            return
        task_name = "综合分析" if intent == KNOWLEDGE_INTENT_GENERAL else "报告检索"
        yield event({"type": "stage", "text": f"已识别为{task_name}任务，正在基于 {len(candidates)} 篇相关报告生成回答…"})
        system, prompt = _knowledge_qa_messages(question, candidates, intent=intent)
        max_tokens = 3600 if intent == KNOWLEDGE_INTENT_GENERAL else 2800
        # 阶段二：流式生成。标记可能跨 chunk 分裂，扣留缓冲保证
        # ===SOURCE_IDS=== 之后的内容不会混入推送给前端的答案正文。
        raw = ""
        emitted = 0
        got_delta = False
        try:
            for delta in _stream_llm(prompt, system=system, max_tokens=max_tokens):
                got_delta = True
                raw += delta
                if KNOWLEDGE_SOURCE_MARKER in raw:
                    safe_end = raw.index(KNOWLEDGE_SOURCE_MARKER)
                else:
                    safe_end = max(len(raw) - len(KNOWLEDGE_SOURCE_MARKER), 0)
                if safe_end > emitted:
                    yield event({"type": "delta", "text": raw[emitted:safe_end]})
                    emitted = safe_end
            if KNOWLEDGE_SOURCE_MARKER not in raw and len(raw) > emitted:
                yield event({"type": "delta", "text": raw[emitted:]})
            result = _parse_knowledge_answer(raw, candidates)
            store.add_qa_history(user_id, question, result.get("answer", ""), result.get("sources", []))
            yield event({"type": "done", **result, "limit": limit, "used": used + 1, "remaining": remaining})
        except Exception as exc:
            # 不把底层 WinError/代理地址直接暴露给用户，便于定位并保持提示可读。
            yield event({"type": "error", "message": f"知识搜索暂时不可用：{exc}"})
        finally:
            # 模型已开始生成（LLM 成本已发生）即计入当日额度，
            # 客户端中途断开也不能通过反复重试绕过限制。
            if got_delta:
                store.add_qa_usage(user_id, day, question)

    response = Response(stream_with_context(generate()), mimetype="text/event-stream")
    response.headers["Cache-Control"] = "no-cache"
    # 反向代理（如 nginx）默认缓冲响应体会导致 SSE 事件攒批到达，需显式关闭。
    response.headers["X-Accel-Buffering"] = "no"
    return response


def _remove_report_file(report):
    if not report.get("fileStored"):
        return
    rel = report.get("fileUrl", "")
    if not rel:
        return
    full = report_file_path(report)
    if full:
        try:
            os.remove(full)
        except OSError:
            pass


# --------------------------------------------------------------------------- #
# 评分 API
# --------------------------------------------------------------------------- #
@app.route("/api/ratings", methods=["GET"])
def api_ratings_list():
    user = require_user()
    if not user:
        return json_error("未登录", 401)

    active_reports = [report for report in store.reports()
                      if report.get("reportType", "internal") == "internal"]
    active_report_ids = {report["id"] for report in active_reports}
    # 回收站报告的评分继续保存在数据库中，但前台暂不展示或计入统计；
    # 恢复报告后这些评分会原样重新出现。
    all_ratings = [rating for rating in store.ratings()
                   if rating.get("reportId") in active_report_ids]
    # 应评分人员总数随报告所选打分部门变化，因此按报告分别统计。
    scorer_ids_global = {u["id"] for u in store.users() if u.get("role") != "admin"}
    report_progress = {}
    # 月报打分已关闭：仅深度报告参与评分
    for report in active_reports:
        if report.get("category") != "deep":
            continue
        eligible_ids = {u["id"] for u in eligible_scorers(report)}
        report_rows = [r for r in all_ratings if r.get("reportId") == report["id"]]
        done = len({r.get("userId") for r in report_rows if r.get("userId") in eligible_ids})
        total_scorers = len(eligible_ids)
        report_progress[report["id"]] = {
            "done": done,
            "pending": max(total_scorers - done, 0),
            "total": total_scorers,
        }

    # 报告作者可查看本人报告的团队总分及各维度均分。这里返回报告级总分
    # 和三维均分，但不下发任何评分人的记录；行政代上传时 authorId 已是被
    # 代传人，因而被代传人同样能获得本人报告的得分。
    authored_report_ids = {
        report["id"]
        for report in active_reports
        if report.get("authorId") == user.get("id")
        and report.get("category") == "deep"
    }
    report_scores = {}
    for report_id in authored_report_ids:
        report_rows = [r for r in all_ratings if r.get("reportId") == report_id]
        if not report_rows:
            report_scores[report_id] = {
                "overall": None,
                "inspiration": None,
                "depth": None,
                "utility": None,
            }
            continue
        overall = sum(
            (int(row["inspiration"]) + int(row["depth"]) + int(row["utility"])) / 3
            for row in report_rows
        ) / len(report_rows)
        dim_avg = lambda key: round(
            sum(int(row[key]) for row in report_rows) / len(report_rows), 1
        )
        report_scores[report_id] = {
            "overall": round(overall, 1),
            "inspiration": dim_avg("inspiration"),
            "depth": dim_avg("depth"),
            "utility": dim_avg("utility"),
        }

    is_leader = user.get("role") == "leader"
    if is_leader:
        users_by_id = {u["id"]: u for u in store.users()}
        visible_ratings = []
        for rating in all_ratings:
            item = dict(rating)
            scorer = users_by_id.get(rating.get("userId"), {})
            item["userName"] = scorer.get("name", "已删除用户")
            item["userOrg"] = scorer.get("org", "")
            visible_ratings.append(item)
        anonymous_feedback = []
    else:
        # 普通成员只需取回自己的评分以支持“已评/修改”；行政不取回任何具体分数。
        visible_ratings = ([r for r in all_ratings if r.get("userId") == user["id"]]
                           if user.get("role") != "admin" else [])
        anonymous_feedback = []
        for rating in all_ratings:
            if rating.get("userId") == user["id"] or not (rating.get("comment") or "").strip():
                continue
            anonymous_feedback.append({
                "reportId": rating.get("reportId"),
                "comment": rating.get("comment", ""),
                "updatedAt": rating.get("updatedAt"),
            })

    return jsonify({
        "ratings": visible_ratings,
        "feedback": anonymous_feedback,
        "reportProgress": report_progress,
        "reportScores": report_scores,
        "summary": {
            "totalRatings": len(all_ratings),
            "participants": len({r.get("userId") for r in all_ratings if r.get("userId") in scorer_ids_global}),
            "ratedReports": len({r.get("reportId") for r in all_ratings}),
            "totalScorers": len(scorer_ids_global),
        },
    })


@app.route("/api/ratings", methods=["POST"])
def api_rating_submit():
    user = require_user()
    if not user:
        return json_error("未登录", 401)
    if user.get("role") == "admin":
        return json_error("行政账号不参与评分", 403)

    data = request.get_json(silent=True) or {}
    report_id = data.get("reportId", "")
    report = store.get_report(report_id)
    if not report:
        return json_error("报告不存在", 404)
    if report.get("reportType", "internal") != "internal" or report.get("category") != "deep":
        return json_error("该报告无需评分", 400)
    # 仅所选打分部门的研究人员（及领导）具备评分资格；报告对所有人可见。
    if not is_eligible_scorer(report, user):
        return json_error("您不在此报告的打分人员范围内", 403)

    try:
        inspiration = int(data.get("inspiration"))
        depth = int(data.get("depth"))
        utility = int(data.get("utility"))
    except (TypeError, ValueError):
        return json_error("评分数据无效", 400)

    for v in (inspiration, depth, utility):
        if not (1 <= v <= 10):
            return json_error("评分需在 1-10 之间", 400)

    comment = (data.get("comment") or "").strip()[:500]
    if store.has_rating(report_id, user["id"]):
        return json_error("该报告您已提交过评分，评分提交后不可修改", 400)
    record = {
        "id": f"rating-{int(time.time()*1000)}-{uuid.uuid4().hex[:6]}",
        "reportId": report_id,
        "userId": user["id"],
        "inspiration": inspiration,
        "depth": depth,
        "utility": utility,
        "comment": comment,
        "updatedAt": _now_iso(),
    }
    store.add_rating(record)
    return jsonify({"ok": True, "rating": record})


@app.route("/api/reports/<rid>/scoring-status", methods=["GET"])
def api_scoring_status(rid):
    """查询某份报告的评分进度：哪些应评分人员已评/未评。
    仅行政与领导可访问（不泄露给普通成员）。
    """
    user = require_role("admin", "leader")
    if not user:
        return json_error("无权限", 403)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    if report.get("reportType", "internal") != "internal":
        return json_error("外部报告不参与评分", 400)
    # 应评分人员按报告所选打分部门确定（部门领导按 org 过滤，通用领导始终参与）
    scorers_all = eligible_scorers(report)
    ratings = [r for r in store.ratings() if r.get("reportId") == rid]
    rating_by_user = {r.get("userId"): r for r in ratings}
    scorers = []
    for u in scorers_all:
        r = rating_by_user.get(u["id"])
        scorers.append({
            "id": u["id"],
            "name": u["name"],
            "org": u["org"],
            "scored": r is not None,
            "updatedAt": r.get("updatedAt") if r else None,
            "score": (round((int(r["inspiration"]) + int(r["depth"]) + int(r["utility"])) / 3, 1)
                      if r else None),
        })
    done = sum(1 for s in scorers if s["scored"])
    response = {
        "reportId": rid,
        "total": len(scorers),
        "done": done,
        "pending": len(scorers) - done,
    }
    # 领导可查看具名评分进度和具体得分；行政可查看具名评分进度但不展示具体得分。
    if user.get("role") == "leader":
        response["scorers"] = scorers
    elif user.get("role") == "admin":
        response["scorers"] = [{**s, "score": None} for s in scorers]
    return jsonify(response)


# --------------------------------------------------------------------------- #
# 行政角色批量操作
# --------------------------------------------------------------------------- #
@app.route("/api/admin/batch-delete", methods=["POST"])
def api_batch_delete():
    user = require_role("admin")
    if not user:
        return json_error("无权限", 403)
    data = request.get_json(silent=True) or {}
    ids = data.get("ids", [])
    if not isinstance(ids, list) or not ids:
        return json_error("未选择报告", 400)
    trashed = store.trash_reports(ids, user.get("id"))
    return jsonify({"ok": True, "trashed": trashed})


@app.route("/api/admin/batch-category", methods=["POST"])
def api_batch_category():
    user = require_role("admin")
    if not user:
        return json_error("无权限", 403)
    data = request.get_json(silent=True) or {}
    ids = data.get("ids", [])
    category = data.get("category", "")
    if not isinstance(ids, list) or not ids:
        return json_error("未选择报告", 400)
    if category not in ("weekly", "monthly", "deep", "other"):
        return json_error("分类无效", 400)
    for rid in ids:
        store.update_report(rid, {"category": category})
    return jsonify({"ok": True})


@app.route("/api/admin/batch-upload-time", methods=["POST"])
def api_batch_upload_time():
    user = require_role("admin")
    if not user:
        return json_error("无权限", 403)
    data = request.get_json(silent=True) or {}
    ids = data.get("ids", [])
    uploaded_at = data.get("uploadedAt", "")
    if not isinstance(ids, list) or not ids:
        return json_error("未选择报告", 400)
    if not uploaded_at:
        return json_error("请提供上传时间", 400)
    # 允许 datetime-local 的 "YYYY-MM-DDTHH:MM" 格式
    for rid in ids:
        store.update_report(rid, {"uploadedAt": uploaded_at})
    return jsonify({"ok": True})


# --------------------------------------------------------------------------- #
# 预览转换 API
# --------------------------------------------------------------------------- #
@app.route("/api/preview", methods=["GET", "POST"])
def api_preview():
    if not require_user():
        return json_error("未登录", 401)
    if not SOFFICE:
        return json_error("服务器未安装 LibreOffice，无法转换文件", 500)

    work_dir = None
    try:
        if request.method == "GET":
            rel_path = request.args.get("file", "")
            if not rel_path:
                return json_error("缺少 file 参数", 400)
            normalized_rel = rel_path.replace("\\", "/").lstrip("/")
            report = next((item for item in store.reports()
                           if str(item.get("fileUrl", "")).replace("\\", "/").lstrip("/") == normalized_rel), None)
            if not report:
                return json_error("报告不存在或已在回收站", 404)
            file_path = report_file_path(report)
            if not file_path:
                return json_error("文件不存在", 404)

            ext = os.path.splitext(file_path)[1].lower()
            if ext == ".pdf":
                return send_file(file_path, mimetype="application/pdf")
            if ext not in CONVERTIBLE_EXTS:
                return json_error(f"不支持的格式: {ext}", 400)

            # 永久缓存命中时直接复用。
            cached_pdf = _get_cached_pdf(file_path)
            if cached_pdf:
                return send_file(cached_pdf, mimetype="application/pdf")

            work_dir = _make_preview_work_dir()
            try:
                pdf_path = convert_to_pdf(file_path, work_dir)
                # 转换成功后落盘到持久缓存（失败不影响本次预览）
                cached_pdf = _store_cached_pdf(pdf_path, file_path, report.get("id"))
                if cached_pdf:
                    # 源 PDF 已 move 到缓存目录，直接发送缓存文件；清理临时工作目录
                    _remove_preview_artifacts(work_dir, retries=1, retry_delay=0)
                    return send_file(cached_pdf, mimetype="application/pdf")
                # 落盘失败则回退到一次性发送
                return _send_temporary_pdf(pdf_path, work_dir)
            except Exception:
                _remove_preview_artifacts(work_dir, retries=1, retry_delay=0)
                raise

        else:
            uploaded = request.files.get("file")
            if not uploaded:
                return json_error("缺少上传文件", 400)

            original_name = uploaded.filename or "upload"
            ext = os.path.splitext(original_name)[1].lower()
            if ext == ".pdf":
                return send_file(uploaded.stream, mimetype="application/pdf",
                                 as_attachment=False)
            if ext not in CONVERTIBLE_EXTS:
                return json_error(f"不支持的格式: {ext}", 400)

            work_dir = _make_preview_work_dir()
            tmp_path = os.path.join(work_dir, "source" + ext)
            uploaded.save(tmp_path)
            report_id = str(request.form.get("reportId", "")).strip() or None
            report = store.get_report(report_id) if report_id else None
            if report and report_file_path(report):
                source_path = report_file_path(report)
                cached_pdf = _get_cached_pdf(source_path)
                if cached_pdf:
                    _remove_preview_artifacts(work_dir, retries=1, retry_delay=0)
                    return send_file(cached_pdf, mimetype="application/pdf")
            pdf_path = convert_to_pdf(tmp_path, work_dir)
            if report and report_file_path(report):
                source_path = report_file_path(report)
                cached_pdf = _store_cached_pdf(pdf_path, source_path, report_id)
                if cached_pdf:
                    _remove_preview_artifacts(work_dir, retries=1, retry_delay=0)
                    return send_file(cached_pdf, mimetype="application/pdf")
            return _send_temporary_pdf(pdf_path, work_dir)

    except subprocess.TimeoutExpired:
        if work_dir:
            _schedule_preview_cleanup(work_dir)
        return json_error("转换超时，文件可能过大", 504)
    except Exception as e:
        if work_dir:
            _schedule_preview_cleanup(work_dir)
        return json_error(str(e), 500)


# --------------------------------------------------------------------------- #
# 超级管理员后台 API
# --------------------------------------------------------------------------- #
@app.route("/api/admin/login", methods=["POST"])
def api_admin_login():
    data = request.get_json(silent=True) or {}
    password = data.get("password", "")
    password_hash = store.admin_password_hash()
    if not password_hash:
        return json_error("后台管理员密码尚未初始化，请先执行账号迁移", 503)
    if not check_password_hash(password_hash, password):
        return json_error("后台密码不正确", 401)
    session["internal_knowledge_base_superadmin"] = True
    return jsonify({"ok": True})


@app.route("/api/admin/logout", methods=["POST"])
def api_admin_logout():
    session.pop("internal_knowledge_base_superadmin", None)
    return jsonify({"ok": True})


@app.route("/api/admin/status", methods=["GET"])
def api_admin_status():
    return jsonify({"admin": bool(session.get("internal_knowledge_base_superadmin"))})


@app.route("/api/admin/pdf-cache", methods=["GET", "DELETE"])
def api_admin_pdf_cache():
    if not require_admin():
        return json_error("未登录后台", 401)
    if request.method == "DELETE":
        keys = store.clear_pdf_caches()
        deleted = _delete_cache_files(keys)
        return jsonify({"ok": True, "deleted": deleted})
    reports = {item["id"]: item for item in store.reports(include_deleted=True)}
    items = []
    total_size = 0
    for row in store.pdf_caches():
        report = reports.get(row.get("report_id")) or {}
        item = {
            "cacheKey": row["cache_key"], "reportId": row.get("report_id"),
            "reportTitle": report.get("title") or "—", "fileName": row["file_name"],
            "generatedAt": row["generated_at"], "lastAccessedAt": row["last_accessed_at"],
            "sizeBytes": row["size_bytes"], "conversionVersion": row["conversion_version"],
        }
        total_size += row["size_bytes"]
        items.append(item)
    return jsonify({"items": items, "count": len(items), "totalSizeBytes": total_size})


@app.route("/api/admin/pdf-cache/<cache_key>", methods=["DELETE"])
def api_admin_pdf_cache_delete(cache_key):
    if not require_admin():
        return json_error("未登录后台", 401)
    if not re.fullmatch(r"[0-9a-f]{64}", cache_key):
        return json_error("缓存标识无效", 400)
    existed = store.delete_pdf_cache(cache_key)
    deleted = _delete_cache_files([cache_key])
    return jsonify({"ok": True, "deleted": bool(existed or deleted)})


@app.route("/api/admin/reports/<rid>/pdf-cache", methods=["DELETE"])
def api_admin_report_pdf_cache_delete(rid):
    if not require_admin():
        return json_error("未登录后台", 401)
    keys = store.delete_pdf_caches_for_report(rid)
    deleted = _delete_cache_files(keys)
    return jsonify({"ok": True, "deleted": deleted})


@app.route("/api/admin/change-password", methods=["POST"])
def api_admin_change_password():
    if not require_admin():
        return json_error("未登录后台", 401)
    data = request.get_json(silent=True) or {}
    old_password = str(data.get("oldPassword", ""))
    new_password = str(data.get("newPassword", ""))
    confirm_password = str(data.get("confirmPassword", ""))
    current_hash = store.admin_password_hash() or ""
    if not check_password_hash(current_hash, old_password):
        return json_error("当前管理员密码不正确", 400)
    if len(new_password) < 8:
        return json_error("新密码至少需要8位", 400)
    if new_password != confirm_password:
        return json_error("两次输入的新密码不一致", 400)
    if new_password == old_password:
        return json_error("新密码不能与当前密码相同", 400)
    store.set_admin_password_hash(generate_password_hash(new_password))
    return jsonify({"ok": True})


@app.route("/api/admin/reports", methods=["GET"])
def api_admin_reports():
    if not require_admin():
        return json_error("未登录后台", 401)
    ratings = store.ratings()
    out = []
    for r in store.reports(include_deleted=True):
        item = public_report(r)
        item["ratingCount"] = sum(1 for x in ratings if x["reportId"] == r["id"])
        out.append(item)
    return jsonify({"reports": out})


@app.route("/api/admin/reports/<rid>", methods=["POST", "PUT"])
def api_admin_report_update(rid):
    if not require_admin():
        return json_error("未登录后台", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    data = request.get_json(silent=True) or {}
    allowed = ["title", "category", "theme", "reportType", "org", "reportDate",
               "uploadedAt", "summary", "recommendation", "tags", "author",
               "sourceAuthor", "sourceInstitution"]
    fields = {}
    for k in allowed:
        if k in data:
            v = data[k]
            if k == "category" and v not in ("weekly", "monthly", "deep", "other"):
                return json_error("分类无效", 400)
            if k == "theme" and v not in REPORT_THEMES:
                return json_error("主题无效", 400)
            if k == "reportType" and v not in ("internal", "external", "research_visit", "roadshow"):
                return json_error("报告类型无效", 400)
            if k == "org" and v not in ("资产配置部", "固收中心"):
                return json_error("部门无效", 400)
            if k == "tags":
                if isinstance(v, str):
                    v = [t.strip() for t in re.split(r"[,，]", v) if t.strip()][:8]
                elif isinstance(v, list):
                    v = [str(t).strip() for t in v if str(t).strip()][:8]
            fields[k] = v
    target_type = fields.get("reportType", report.get("reportType", "internal"))
    if target_type == "external":
        if not str(fields.get("sourceAuthor", report.get("sourceAuthor", ""))).strip():
            return json_error("外部报告必须填写原作者", 400)
        if not str(fields.get("sourceInstitution", report.get("sourceInstitution", ""))).strip():
            return json_error("外部报告必须填写机构", 400)
    updated = store.update_report(rid, fields)
    return jsonify({"ok": True, "report": public_report(updated)})


@app.route("/api/admin/reports/<rid>", methods=["DELETE"])
def api_admin_report_delete(rid):
    if not require_admin():
        return json_error("未登录后台", 401)
    report = store.get_report(rid, include_deleted=True)
    if not report:
        return json_error("报告不存在", 404)
    if not report.get("deletedAt"):
        return json_error("请先将报告移入回收站，再确认永久删除", 400)
    _remove_report_file(report)
    _delete_cache_files(store.delete_pdf_caches_for_report(rid))
    store.delete_reports([rid])
    return jsonify({"ok": True, "permanent": True})


@app.route("/api/admin/reports/<rid>/trash", methods=["POST"])
def api_admin_report_trash(rid):
    if not require_admin():
        return json_error("未登录后台", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在或已在回收站", 404)
    store.trash_reports([rid], "superadmin")
    return jsonify({"ok": True, "trashed": True})


@app.route("/api/admin/reports/<rid>/restore", methods=["POST"])
def api_admin_report_restore(rid):
    if not require_admin():
        return json_error("未登录后台", 401)
    report = store.get_report(rid, include_deleted=True)
    if not report:
        return json_error("报告不存在", 404)
    if not report.get("deletedAt"):
        return json_error("报告不在回收站", 400)
    restored = store.restore_report(rid)
    return jsonify({"ok": True, "report": public_report(restored)})


@app.route("/api/admin/reports/<rid>/reset-scores", methods=["POST"])
def api_admin_reset_scores(rid):
    if not require_admin():
        return json_error("未登录后台", 401)
    report = store.get_report(rid)
    if not report:
        return json_error("报告不存在", 404)
    store.reset_ratings_for_report(rid)
    return jsonify({"ok": True})


@app.route("/api/admin/users", methods=["GET"])
def api_admin_users():
    if not require_admin():
        return json_error("未登录后台", 401)
    users = store.users()
    out = []
    for u in users:
        out.append({
            "id": u["id"], "name": u["name"], "org": u["org"], "role": u["role"],
            "defaultPassword": check_password_hash(u.get("password_hash", ""), DEFAULT_USER_PASSWORD),
        })
    return jsonify({"users": out})


@app.route("/api/admin/users", methods=["POST"])
def api_admin_user_create():
    if not require_admin():
        return json_error("未登录后台", 401)
    data = request.get_json(silent=True) or {}
    uid = (data.get("id") or "").strip().lower().replace(" ", "")
    name = (data.get("name") or "").strip()
    org = (data.get("org") or "").strip()
    role = (data.get("role") or "").strip()
    password = (data.get("password") or "").strip() or DEFAULT_USER_PASSWORD
    if not uid or not name:
        return json_error("账号和姓名必填", 400)
    if not re.match(r"^[a-z0-9_.-]+$", uid):
        return json_error("账号只能包含小写字母、数字及 _ . -", 400)
    if org not in ("资产配置部", "固收中心", "领导", "行政", ""):
        return json_error("部门无效", 400)
    if role not in ("leader", "admin", "member"):
        return json_error("角色无效", 400)
    if store.get_user(uid):
        return json_error("账号已存在", 400)
    user = {"id": uid, "name": name, "org": org, "role": role,
            "password_hash": generate_password_hash(password)}
    store.add_user(user)
    return jsonify({"ok": True, "user": public_user(user)})


@app.route("/api/admin/users/<uid>", methods=["POST", "PUT"])
def api_admin_user_update(uid):
    if not require_admin():
        return json_error("未登录后台", 401)
    user = store.get_user(uid)
    if not user:
        return json_error("用户不存在", 404)
    data = request.get_json(silent=True) or {}
    fields = {}
    for k in ("name", "org", "role"):
        if k in data:
            v = (data[k] or "").strip() if isinstance(data[k], str) else data[k]
            if k == "org" and v not in ("资产配置部", "固收中心", "领导", "行政", ""):
                return json_error("部门无效", 400)
            if k == "role" and v not in ("leader", "admin", "member"):
                return json_error("角色无效", 400)
            if k == "name" and not v:
                return json_error("姓名不能为空", 400)
            fields[k] = v
    updated = store.update_user(uid, fields)
    return jsonify({"ok": True, "user": public_user(updated)})


@app.route("/api/admin/users/<uid>", methods=["DELETE"])
def api_admin_user_delete(uid):
    if not require_admin():
        return json_error("未登录后台", 401)
    if not store.get_user(uid):
        return json_error("用户不存在", 404)
    store.delete_user(uid)
    return jsonify({"ok": True})


@app.route("/api/admin/users/<uid>/reset-password", methods=["POST"])
def api_admin_user_reset_password(uid):
    if not require_admin():
        return json_error("未登录后台", 401)
    user = store.get_user(uid)
    if not user:
        return json_error("用户不存在", 404)
    data = request.get_json(silent=True) or {}
    password = (data.get("password") or "").strip()
    if not password:
        return json_error("请输入新密码", 400)
    if len(password) < 6:
        return json_error("密码至少 6 位", 400)
    store.update_user(uid, {"password_hash": generate_password_hash(password)})
    return jsonify({"ok": True})


@app.route("/api/admin/users/<uid>/clear-qa-history", methods=["POST"])
def api_admin_user_clear_qa_history(uid):
    if not require_admin():
        return json_error("未登录后台", 401)
    if not store.get_user(uid):
        return json_error("用户不存在", 404)
    removed = store.clear_qa_history(uid)
    return jsonify({"ok": True, "removed": removed})


@app.route("/api/admin/reminder-config", methods=["GET", "PUT", "POST"])
def api_admin_reminder_config():
    if not require_admin():
        return json_error("未登录后台", 401)
    if request.method == "GET":
        return jsonify({"config": store.reminder_config()})
    data = request.get_json(silent=True) or {}
    period = str(data.get("period", "")).strip()
    if not re.match(r"^\d{4}$", period):
        return json_error("统计年度格式应为 YYYY", 400)
    category = str(data.get("reportCategory", "deep")).strip()
    if category not in ("weekly", "monthly", "deep", "other"):
        return json_error("专题报告分类无效", 400)
    valid_user_ids = {item["id"] for item in store.users()}
    clean_rules = []
    rules = data.get("rules", [])
    if not isinstance(rules, list):
        return json_error("规则格式无效", 400)
    for index, rule in enumerate(rules[:100]):
        if not isinstance(rule, dict):
            continue
        label = str(rule.get("label", "")).strip()[:80]
        if not label:
            return json_error(f"第 {index + 1} 条规则缺少名称", 400)
        try:
            target = int(rule.get("target", 0))
        except (TypeError, ValueError):
            return json_error(f"第 {index + 1} 条规则目标数量无效", 400)
        if not 0 <= target <= 100:
            return json_error("目标数量应在 0-100 之间", 400)
        user_ids = list(dict.fromkeys(
            str(uid) for uid in (rule.get("userIds") or []) if str(uid) in valid_user_ids
        ))
        clean_rules.append({
            "id": str(rule.get("id") or f"rule-{uuid.uuid4().hex[:8]}"),
            "label": label, "mode": "person" if rule.get("mode") == "person" else "group",
            "target": target, "userIds": user_ids,
        })
    config = {"period": period, "reportCategory": category, "rules": clean_rules}
    store.set_reminder_config(config)
    return jsonify({"ok": True, "config": config})


@app.route("/api/admin/knowledge-config", methods=["GET", "PUT", "POST"])
def api_admin_knowledge_config():
    if not require_admin():
        return json_error("未登录后台", 401)
    if request.method == "GET":
        return jsonify({"config": store.knowledge_config()})
    data = request.get_json(silent=True) or {}
    try:
        member_limit = int(data.get("memberLimit", 10))
        leader_limit = int(data.get("leaderLimit", 100))
    except (TypeError, ValueError):
        return json_error("知识搜索额度必须是整数", 400)
    if not 1 <= member_limit <= 1000 or not 1 <= leader_limit <= 1000:
        return json_error("知识搜索额度应在 1 至 1000 次之间", 400)
    config = {"memberLimit": member_limit, "leaderLimit": leader_limit}
    store.set_knowledge_config(config)
    return jsonify({"ok": True, "config": config})


@app.route("/api/admin/stats", methods=["GET"])
def api_admin_stats():
    if not require_admin():
        return json_error("未登录后台", 401)
    reports = store.reports()
    all_reports = store.reports(include_deleted=True)
    users = store.users()
    active_ids = {report["id"] for report in reports}
    ratings = [rating for rating in store.ratings() if rating.get("reportId") in active_ids]
    return jsonify({
        "reports": len(reports),
        "internalReports": sum(1 for report in reports if report.get("reportType", "internal") == "internal"),
        "externalReports": sum(1 for report in reports if report.get("reportType") == "external"),
        "deletedReports": sum(1 for report in all_reports if report.get("deletedAt")),
        "users": len(users),
        "ratings": len(ratings),
        "ratedReports": len({r["reportId"] for r in ratings}),
        "participants": len({r["userId"] for r in ratings}),
        "byCategory": {c: sum(1 for r in reports if r.get("category") == c)
                       for c in ("weekly", "monthly", "deep", "other")},
        "byTheme": {t: sum(1 for r in reports if r.get("theme") == t)
                    for t in REPORT_THEMES},
    })
